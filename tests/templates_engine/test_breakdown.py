# tests/templates_engine/test_breakdown.py
# Phase D+E+F: Tests for pptx structural inspector, content stripping,
#              usability scoring, draft manifest, detection report, and hard stop.
from pathlib import Path
from unittest.mock import patch

import pptx.presentation
import pytest
import yaml
from pptx import Presentation

from app.services.file_manipulator import ManipulationError
from templates_engine.breakdown import (
    detection_report,
    format_inspection,
    generate_draft_manifest,
    hard_stop_result,
    inspect_pptx,
    score_layouts,
    strip_to_template,
)
from templates_engine.manifest import load_manifest

# Bundled default template shipped with Simplicitor -- a real .pptx fixture.
BUNDLED_TEMPLATE = (
    Path(__file__).parent.parent.parent / "simplicitor" / "templates" / "pptx_default.pptx"
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_minimal_pptx(path: Path) -> Path:
    """Save a default-layout Presentation to path and return it."""
    prs = Presentation()  # uses python-pptx's built-in blank template
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# inspect_pptx — happy path
# ---------------------------------------------------------------------------

def test_inspect_returns_dict(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    assert isinstance(report, dict)
    assert "path" in report
    assert "layouts" in report


def test_inspect_has_layouts(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    assert len(report["layouts"]) > 0


def test_inspect_layout_has_required_keys(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    layout = report["layouts"][0]
    assert "layout_index" in layout
    assert "name" in layout
    assert "placeholders" in layout


def test_inspect_layout_index_is_sequential(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    for i, layout in enumerate(report["layouts"]):
        assert layout["layout_index"] == i


def test_inspect_placeholder_has_required_keys(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    # Find any layout that has placeholders
    ph = next(
        ph
        for layout in report["layouts"]
        for ph in layout["placeholders"]
    )
    assert "idx" in ph
    assert "type" in ph
    assert "name" in ph
    assert "position" in ph
    assert "is_custom" in ph


def test_inspect_detects_title_at_idx_0(tmp_path):
    """Layout 0 of a default presentation must have a placeholder with idx=0."""
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    layout_0 = report["layouts"][0]
    idx_values = [ph["idx"] for ph in layout_0["placeholders"]]
    assert 0 in idx_values, (
        f"Expected placeholder idx=0 in layout 0. Got indices: {idx_values}"
    )


def test_inspect_idx_0_type_contains_title(tmp_path):
    """The idx=0 placeholder in layout 0 should have a title-related type."""
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    layout_0 = report["layouts"][0]
    ph0 = next(ph for ph in layout_0["placeholders"] if ph["idx"] == 0)
    assert "TITLE" in ph0["type"].upper(), (
        f"Expected TITLE in type of idx=0 placeholder, got: {ph0['type']!r}"
    )


def test_inspect_is_custom_true_for_idx_gte_10(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    for layout in report["layouts"]:
        for ph in layout["placeholders"]:
            expected = ph["idx"] >= 10
            assert ph["is_custom"] == expected, (
                f"is_custom mismatch for idx={ph['idx']}: expected {expected}"
            )


def test_inspect_position_dict_has_four_keys(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    ph = next(
        ph for layout in report["layouts"] for ph in layout["placeholders"]
    )
    pos = ph["position"]
    assert set(pos.keys()) == {"left", "top", "width", "height"}


def test_inspect_path_is_absolute(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    assert Path(report["path"]).is_absolute()


def test_inspect_bundled_template_smoke():
    """Smoke test against the real bundled template shipped with Simplicitor."""
    assert BUNDLED_TEMPLATE.exists(), f"Bundled template not found at {BUNDLED_TEMPLATE}"
    report = inspect_pptx(BUNDLED_TEMPLATE)
    assert len(report["layouts"]) >= 1
    # Layout 0 must have a title placeholder at idx 0
    layout_0 = report["layouts"][0]
    assert any(ph["idx"] == 0 for ph in layout_0["placeholders"])


# ---------------------------------------------------------------------------
# inspect_pptx — error paths (conventional Simplicitor errors)
# ---------------------------------------------------------------------------

def test_missing_file_raises_manipulation_error(tmp_path):
    with pytest.raises(ManipulationError, match=r"not found|missing"):
        inspect_pptx(tmp_path / "ghost.pptx")


def test_non_pptx_extension_raises_value_error(tmp_path):
    txt = tmp_path / "notes.txt"
    txt.write_text("not a presentation", encoding="utf-8")
    with pytest.raises(ValueError, match=r"\.pptx|extension"):
        inspect_pptx(txt)


def test_non_pptx_error_names_extension(tmp_path):
    pdf = tmp_path / "report.pdf"
    pdf.write_bytes(b"%PDF-1.4 fake content")
    with pytest.raises(ValueError, match=r"\.pdf|extension|Expected"):
        inspect_pptx(pdf)


def test_corrupt_pptx_raises_value_error(tmp_path):
    bad = tmp_path / "corrupt.pptx"
    bad.write_bytes(b"this is not a zip file at all")
    with pytest.raises(ValueError, match=r"[Cc]ould not open|PowerPoint"):
        inspect_pptx(bad)


def test_corrupt_pptx_valid_zip_wrong_structure(tmp_path):
    """A valid zip that is not a pptx internally must also raise ValueError."""
    import zipfile as zf
    bad = tmp_path / "fake.pptx"
    with zf.ZipFile(str(bad), "w") as z:
        z.writestr("not_a_pptx.txt", "wrong content")
    with pytest.raises(ValueError, match=r"[Cc]ould not open|PowerPoint"):
        inspect_pptx(bad)


# ---------------------------------------------------------------------------
# format_inspection
# ---------------------------------------------------------------------------

def test_format_inspection_is_string(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    result = format_inspection(report)
    assert isinstance(result, str)


def test_format_inspection_contains_path(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    result = format_inspection(report)
    assert "test.pptx" in result


def test_format_inspection_contains_layout_index(tmp_path):
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    result = format_inspection(report)
    assert "[" in result and "]" in result  # layout index brackets


def test_format_inspection_marks_custom_placeholders(tmp_path):
    """idx >= 10 placeholders should be flagged as CUSTOM."""
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    report = inspect_pptx(pptx)
    # Only check if any custom placeholder exists in the template
    has_custom = any(
        ph["is_custom"]
        for layout in report["layouts"]
        for ph in layout["placeholders"]
    )
    if has_custom:
        result = format_inspection(report)
        assert "CUSTOM" in result


# ---------------------------------------------------------------------------
# Phase E: strip_to_template
# ---------------------------------------------------------------------------

def test_strip_yields_zero_slides(tmp_path):
    source = _make_minimal_pptx(tmp_path / "source.pptx")
    out = tmp_path / "stripped.pptx"
    strip_to_template(source, out)
    result = Presentation(str(out))
    assert len(result.slides) == 0


def test_strip_preserves_slide_masters(tmp_path):
    """Stripping must keep slide masters intact (the 'design-only' guarantee)."""
    source = _make_minimal_pptx(tmp_path / "source.pptx")
    master_count = len(Presentation(str(source)).slide_masters)
    out = tmp_path / "stripped.pptx"
    strip_to_template(source, out)
    assert len(Presentation(str(out)).slide_masters) == master_count


def test_strip_real_slides_yields_zero(tmp_path):
    """Strip a presentation with real slides -- confirms they are actually removed."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)
    source = tmp_path / "two_slides.pptx"
    prs.save(str(source))
    out = tmp_path / "stripped.pptx"
    strip_to_template(source, out)
    result = Presentation(str(out))
    assert len(result.slides) == 0
    assert len(result.slide_masters) > 0


def test_strip_preserves_layouts(tmp_path):
    source = _make_minimal_pptx(tmp_path / "source.pptx")
    original_layout_count = len(Presentation(str(source)).slide_layouts)
    out = tmp_path / "stripped.pptx"
    strip_to_template(source, out)
    stripped = Presentation(str(out))
    assert len(stripped.slide_layouts) == original_layout_count


def test_strip_wrong_out_extension_raises_value_error(tmp_path):
    source = _make_minimal_pptx(tmp_path / "source.pptx")
    with pytest.raises(ValueError, match=r"\.pptx|extension"):
        strip_to_template(source, tmp_path / "output.docx")


def test_strip_write_failure_raises_manipulation_error(tmp_path):
    source = _make_minimal_pptx(tmp_path / "source.pptx")
    out = tmp_path / "output.pptx"
    with patch.object(pptx.presentation.Presentation, "save", side_effect=OSError("simulated disk full")):
        with pytest.raises(ManipulationError):
            strip_to_template(source, out)


def test_strip_write_failure_leaves_no_debris(tmp_path):
    source = _make_minimal_pptx(tmp_path / "source.pptx")
    out = tmp_path / "output.pptx"
    out.write_bytes(b"partial content")  # pre-create to simulate a partial write
    with patch.object(pptx.presentation.Presentation, "save", side_effect=OSError("simulated disk full")):
        with pytest.raises(ManipulationError):
            strip_to_template(source, out)
    assert not out.exists(), "Partial output file must be deleted on write failure"


# ---------------------------------------------------------------------------
# Phase E: score_layouts
# ---------------------------------------------------------------------------

def _fake_inspection(layouts_spec: list[list[dict]]) -> dict:
    """Build a minimal inspection dict from a list of placeholder-spec lists."""
    layouts = []
    for i, placeholders in enumerate(layouts_spec):
        layouts.append({
            "layout_index": i,
            "name": f"Layout {i}",
            "placeholders": [
                {
                    "idx": ph["idx"],
                    "type": ph.get("type", f"BODY ({ph['idx']})"),
                    "name": ph.get("name", f"Placeholder {ph['idx']}"),
                    "position": {"left": 0, "top": 0, "width": 0, "height": 0},
                    "is_custom": ph["idx"] >= 10,
                }
                for ph in placeholders
            ],
        })
    return {"path": "/fake/path.pptx", "layouts": layouts}


def test_score_wrong_input_raises_value_error():
    with pytest.raises(ValueError, match=r"inspect_pptx"):
        score_layouts({"not_layouts": []})


def test_score_title_only_is_unusable():
    inspection = _fake_inspection([[{"idx": 0, "type": "TITLE (1)"}]])
    result = score_layouts(inspection)
    assert result["layouts"][0]["usable"] is False


def test_score_title_and_body_is_usable():
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    result = score_layouts(inspection)
    assert result["layouts"][0]["usable"] is True


def test_score_custom_only_is_unusable():
    inspection = _fake_inspection([[{"idx": 10, "type": "OBJECT (14)"}]])
    result = score_layouts(inspection)
    assert result["layouts"][0]["usable"] is False


def test_score_empty_layout_is_unusable():
    inspection = _fake_inspection([[]])
    result = score_layouts(inspection)
    assert result["layouts"][0]["usable"] is False


def test_score_deck_is_usable_if_any_layout_usable():
    # First layout: title-only (unusable), second: title+body (usable)
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}],
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}],
    ])
    result = score_layouts(inspection)
    assert result["is_usable"] is True


def test_score_deck_is_unusable_when_all_layouts_unusable():
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}],
        [{"idx": 10, "type": "OBJECT (14)"}],
        [],
    ])
    result = score_layouts(inspection)
    assert result["is_usable"] is False


def test_score_layout_with_subtitle_is_usable():
    """SUBTITLE (idx=1) qualifies as a content placeholder."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "SUBTITLE (4)"}]
    ])
    result = score_layouts(inspection)
    assert result["layouts"][0]["usable"] is True


def test_score_decorative_placeholder_not_counted():
    """DATE/FOOTER/SLIDE_NUMBER do not make a layout usable."""
    inspection = _fake_inspection([
        [
            {"idx": 0, "type": "TITLE (1)"},
            {"idx": 2, "type": "DATE (10)"},
            {"idx": 3, "type": "FOOTER (11)"},
            {"idx": 4, "type": "SLIDE_NUMBER (12)"},
        ]
    ])
    result = score_layouts(inspection)
    assert result["layouts"][0]["usable"] is False


def test_score_returns_per_layout_entries_for_all_layouts():
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}],
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}],
    ])
    result = score_layouts(inspection)
    assert len(result["layouts"]) == 2
    assert result["layouts"][0]["layout_index"] == 0
    assert result["layouts"][1]["layout_index"] == 1


def test_score_real_pptx(tmp_path):
    """score_layouts on a real python-pptx generated file produces a valid result."""
    pptx = _make_minimal_pptx(tmp_path / "test.pptx")
    inspection = inspect_pptx(pptx)
    result = score_layouts(inspection)
    assert "is_usable" in result
    assert "layouts" in result
    assert len(result["layouts"]) == len(inspection["layouts"])


# ---------------------------------------------------------------------------
# Phase F: generate_draft_manifest
# ---------------------------------------------------------------------------

def test_draft_manifest_title_auto_labeled():
    """idx=0 placeholder must be named 'title' with kind='text'."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    title_field = next(f for f in fields if f["placeholder_idx"] == 0)
    assert title_field["name"] == "title"
    assert title_field["kind"] == "text"


def test_draft_manifest_body_auto_labeled():
    """Single BODY placeholder must be named 'body' with kind='bullets'."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    body_field = next(f for f in fields if f["placeholder_idx"] == 1)
    assert body_field["name"] == "body"
    assert body_field["kind"] == "bullets"


def test_draft_manifest_picture_auto_labeled():
    """PICTURE placeholder must be named 'image' with kind='image'."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "PICTURE (18)"}]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    img_field = next(f for f in fields if f["placeholder_idx"] == 1)
    assert img_field["name"] == "image"
    assert img_field["kind"] == "image"


def test_draft_manifest_same_type_both_get_needs_label():
    """Two BODY placeholders in the same layout both get NEEDS_LABEL_<idx>."""
    inspection = _fake_inspection([
        [
            {"idx": 0, "type": "TITLE (1)"},
            {"idx": 1, "type": "BODY (2)"},
            {"idx": 2, "type": "BODY (2)"},
        ]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    ambiguous = [f for f in fields if f["placeholder_idx"] in (1, 2)]
    assert all(f["name"].startswith("NEEDS_LABEL_") for f in ambiguous)


def test_draft_manifest_two_picture_placeholders_get_needs_label(tmp_path):
    """Two PICTURE placeholders must both become NEEDS_LABEL, not duplicate 'image' names.

    Duplicate 'image' names would break the load_manifest round-trip.
    """
    inspection = _fake_inspection([
        [
            {"idx": 0, "type": "TITLE (1)"},
            {"idx": 1, "type": "PICTURE (18)"},
            {"idx": 2, "type": "PICTURE (18)"},
        ]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    picture_fields = [f for f in fields if f["placeholder_idx"] in (1, 2)]
    assert all(f["name"].startswith("NEEDS_LABEL_") for f in picture_fields)
    # Verify round-trip: no duplicate names means load_manifest must succeed.
    manifest_path = tmp_path / "manifest.yaml"
    manifest_path.write_text(yaml.dump(result), encoding="utf-8")
    loaded = load_manifest(manifest_path)
    assert len(loaded.slide_types) == 1


def test_draft_manifest_excludes_unusable_layouts():
    """Unusable layouts (title-only) must not appear in slide_types."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}],  # unusable
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}],  # usable
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    assert len(result["slide_types"]) == 1
    slide_type = list(result["slide_types"].values())[0]
    assert slide_type["layout_index"] == 1


def test_draft_manifest_excludes_decorative_placeholders():
    """DATE/FOOTER/SLIDE_NUMBER placeholders must not appear as fields."""
    inspection = _fake_inspection([
        [
            {"idx": 0, "type": "TITLE (1)"},
            {"idx": 1, "type": "BODY (2)"},
            {"idx": 2, "type": "DATE (10)"},
            {"idx": 3, "type": "FOOTER (11)"},
        ]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    field_indices = {f["placeholder_idx"] for f in fields}
    assert 2 not in field_indices
    assert 3 not in field_indices


def test_draft_manifest_excludes_custom_placeholders():
    """Custom placeholders (idx>=10) must not appear as fields."""
    inspection = _fake_inspection([
        [
            {"idx": 0, "type": "TITLE (1)"},
            {"idx": 1, "type": "BODY (2)"},
            {"idx": 11, "type": "OBJECT (14)"},
        ]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    fields = list(result["slide_types"].values())[0]["fields"]
    assert 11 not in {f["placeholder_idx"] for f in fields}


def test_draft_manifest_has_required_top_level_keys():
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    assert result["type"] == "pptx"
    assert result["template_file"] == "test.pptx"
    assert "name" in result
    assert "description" in result
    assert "slide_types" in result


def test_draft_manifest_deduplicated_slide_type_keys():
    """Two layouts with the same name get distinct slide_type keys."""
    inspection = {
        "path": "/fake/path.pptx",
        "layouts": [
            {
                "layout_index": 0,
                "name": "Content",
                "placeholders": [
                    {"idx": 0, "type": "TITLE (1)", "name": "T", "position": {}, "is_custom": False},
                    {"idx": 1, "type": "BODY (2)", "name": "B", "position": {}, "is_custom": False},
                ],
            },
            {
                "layout_index": 1,
                "name": "Content",
                "placeholders": [
                    {"idx": 0, "type": "TITLE (1)", "name": "T", "position": {}, "is_custom": False},
                    {"idx": 1, "type": "BODY (2)", "name": "B", "position": {}, "is_custom": False},
                ],
            },
        ],
    }
    scoring = score_layouts(inspection)
    result = generate_draft_manifest(inspection, scoring, "test.pptx")
    keys = list(result["slide_types"].keys())
    assert len(keys) == 2
    assert len(set(keys)) == 2


def test_draft_manifest_round_trips_via_load_manifest(tmp_path):
    """Draft manifest dict must round-trip through YAML and load_manifest."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    scoring = score_layouts(inspection)
    manifest_dict = generate_draft_manifest(inspection, scoring, "test.pptx")

    manifest_path = tmp_path / "manifest.yaml"
    manifest_path.write_text(yaml.dump(manifest_dict), encoding="utf-8")

    loaded = load_manifest(manifest_path)
    assert loaded.name == manifest_dict["name"]
    assert loaded.type == "pptx"
    assert len(loaded.slide_types) == 1


# ---------------------------------------------------------------------------
# Phase F: detection_report
# ---------------------------------------------------------------------------

def test_detection_report_is_string():
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    scoring = score_layouts(inspection)
    assert isinstance(detection_report(inspection, scoring), str)


def test_detection_report_usable_deck_says_can():
    """Usable deck must have a CAN verdict in the report."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}]
    ])
    scoring = score_layouts(inspection)
    report = detection_report(inspection, scoring)
    assert "CAN" in report


def test_detection_report_unusable_deck_says_cannot():
    """Unusable deck (title-only layout) must have a CANNOT verdict."""
    inspection = _fake_inspection([[{"idx": 0, "type": "TITLE (1)"}]])
    scoring = score_layouts(inspection)
    report = detection_report(inspection, scoring)
    assert "CANNOT" in report


def test_detection_report_contains_layout_counts():
    """Report must mention the total number of layouts."""
    inspection = _fake_inspection([
        [{"idx": 0, "type": "TITLE (1)"}],
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}],
        [{"idx": 0, "type": "TITLE (1)"}, {"idx": 1, "type": "BODY (2)"}],
    ])
    scoring = score_layouts(inspection)
    report = detection_report(inspection, scoring)
    assert "3" in report  # total layouts


def test_detection_report_mentions_needs_label():
    """Report must include NEEDS_LABEL field names when ambiguous placeholders are present."""
    inspection = _fake_inspection([
        [
            {"idx": 0, "type": "TITLE (1)"},
            {"idx": 1, "type": "BODY (2)"},
            {"idx": 2, "type": "BODY (2)"},
        ]
    ])
    scoring = score_layouts(inspection)
    report = detection_report(inspection, scoring)
    assert "NEEDS_LABEL" in report


# ---------------------------------------------------------------------------
# Phase F: hard_stop_result
# ---------------------------------------------------------------------------

def test_hard_stop_result_returns_dict():
    assert isinstance(hard_stop_result(), dict)


def test_hard_stop_result_has_hard_stop_status():
    assert hard_stop_result()["status"] == "hard_stop"


def test_hard_stop_result_contains_verbatim_phrases():
    message = hard_stop_result()["message"]
    assert "can't be used as a template" in message
    assert "none of this file's layouts offer a content area" in message
    assert "layouts include body or content areas" in message


def test_hard_stop_result_is_returned_not_raised():
    """hard_stop_result must return normally, never raise."""
    result = hard_stop_result()
    assert result is not None
