# tests/templates_engine/test_render_pptx.py
# Phase H: Tests for the PPTX renderer.
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation

from app.services.file_manipulator import ManipulationError
from templates_engine.manifest import FieldDef, Manifest, SlideTypeDef, load_manifest
from templates_engine.render_pptx import render


# ---------------------------------------------------------------------------
# Happy path
# ---------------------------------------------------------------------------

def test_render_returns_result_dict(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hello"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert isinstance(result, dict)
    assert "path" in result
    assert "issues" in result


def test_render_output_file_exists(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hello"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert result["path"].is_file()


def test_render_slide_count_matches_content(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {
        "slides": [
            {"type": "title_slide", "fields": {"title": "T1"}},
            {"type": "content_slide", "fields": {"heading": "H1", "body": ["A", "B"]}},
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    prs = Presentation(str(result["path"]))
    assert len(prs.slides) == 2


def test_render_text_field_written_to_slide(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "My Title"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[0].text == "My Title"


def test_render_bullets_written_to_slide(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {
        "slides": [{"type": "content_slide", "fields": {"heading": "H", "body": ["X", "Y"]}}]
    }
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    prs = Presentation(str(result["path"]))
    texts = [p.text for p in prs.slides[0].placeholders[1].text_frame.paragraphs if p.text]
    assert "X" in texts
    assert "Y" in texts


def test_render_no_issues_on_clean_content(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert result["issues"] == []


def test_render_appends_pptx_extension_when_absent(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}
    result = render(manifest, content, tmp_path / "out", tmp_template)
    assert result["path"].suffix == ".pptx"
    assert result["path"].is_file()


def test_render_creates_output_directory_if_needed(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}
    nested = tmp_path / "sub" / "dir" / "out.pptx"
    result = render(manifest, content, nested, tmp_template)
    assert result["path"].is_file()


# ---------------------------------------------------------------------------
# Degrade cases — warn + collect, never raise
# ---------------------------------------------------------------------------

def test_render_missing_optional_field_skipped(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    # subtitle is optional; body is optional — omit both
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert result["path"].is_file()
    assert result["issues"] == []


def test_render_text_overflow_warns_not_truncates(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    long_title = "X" * 100  # render_manifest has max_chars=20 for title_slide.title
    content = {"slides": [{"type": "title_slide", "fields": {"title": long_title}}]}

    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)

    assert len(result["issues"]) == 1
    assert "max_chars" in result["issues"][0]
    # Full text must be in the rendered file — no truncation
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[0].text == long_title


def test_render_bullets_overflow_warns_not_caps(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    bullets = ["A", "B", "C", "D"]  # render_manifest has max_items=2 for content_slide.body
    content = {
        "slides": [{"type": "content_slide", "fields": {"heading": "H", "body": bullets}}]
    }

    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)

    assert len(result["issues"]) == 1
    assert "max_items" in result["issues"][0]
    # All 4 bullets must be in the rendered file — no cap
    prs = Presentation(str(result["path"]))
    texts = [p.text for p in prs.slides[0].placeholders[1].text_frame.paragraphs if p.text]
    assert len(texts) == 4
    assert texts == ["A", "B", "C", "D"]


def test_render_multiple_degrade_warnings_collected(tmp_path, tmp_template):
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {
        "slides": [
            {"type": "title_slide", "fields": {"title": "X" * 100}},
            {"type": "content_slide", "fields": {"heading": "H", "body": ["A", "B", "C", "D"]}},
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert len(result["issues"]) == 2
    assert any("max_chars" in i for i in result["issues"])
    assert any("max_items" in i for i in result["issues"])


# ---------------------------------------------------------------------------
# Image field tests
# ---------------------------------------------------------------------------

def test_render_missing_image_warns_not_raises(tmp_path, tmp_template):
    """Image path does not exist — warn + skip, no exception, slide still rendered."""
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {
        "slides": [
            {"type": "image_slide", "fields": {"heading": "Hi", "img": str(tmp_path / "no_such_image.png")}}
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert result["path"].is_file()
    assert len(result["issues"]) == 1
    assert "not found" in result["issues"][0]
    # Heading was rendered despite missing image — degrade affected only the image field
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[0].text == "Hi"


def test_render_image_insert_failure_warns_not_raises(tmp_path, tmp_template, tiny_png):
    """insert_picture raises on a non-picture placeholder — warn + skip, no exception.

    No mock needed: layout 1's body placeholder at idx=1 has no insert_picture method;
    it raises AttributeError naturally, which the renderer's except-Exception handler
    catches and converts to a degrade warning. The heading field is still rendered.
    """
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {
        "slides": [{"type": "image_slide", "fields": {"heading": "Hi", "img": str(tiny_png)}}]
    }
    result = render(manifest, content, tmp_path / "out.pptx", tmp_template)
    assert result["path"].is_file()
    assert len(result["issues"]) >= 1
    assert any("insert" in i.lower() or "image" in i.lower() or "AttributeError" in i for i in result["issues"])
    # Heading was rendered despite image insert failure — degrade affected only the image field
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[0].text == "Hi"


def test_render_image_field_does_not_raise(tmp_path, tmp_template, tiny_png):
    """insert_picture called with the image path; returned placeholder captured; no issues.

    Mocking rationale: render_manifest.yaml maps image_slide to layout_index=1 ("Title and
    Content") which has no PICTURE-type placeholder at idx=1. A narrow mock substitutes only
    the placeholder for the img field (identified by field_name=="img"), preserving real
    placeholder lookup for the heading field (idx=0). This is the minimum mock needed to
    exercise the insert_picture success path given the fixture template constraint.
    """
    import templates_engine.render_pptx as _rp
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {
        "slides": [{"type": "image_slide", "fields": {"heading": "Hi", "img": str(tiny_png)}}]
    }

    real_get_placeholder = _rp._get_placeholder
    mock_img_ph = MagicMock()
    mock_img_ph.insert_picture.return_value = MagicMock()

    def _narrow_get_placeholder(slide, idx, field_name, slide_idx):
        if field_name == "img":
            return mock_img_ph
        return real_get_placeholder(slide, idx, field_name, slide_idx)

    with patch("templates_engine.render_pptx._get_placeholder", side_effect=_narrow_get_placeholder):
        result = render(manifest, content, tmp_path / "out.pptx", tmp_template)

    assert result["issues"] == []
    mock_img_ph.insert_picture.assert_called_once_with(str(tiny_png))


# ---------------------------------------------------------------------------
# Genuine faults — raise, no partial file
# ---------------------------------------------------------------------------

def test_render_bad_layout_index_raises_manipulation_error(tmp_path, tmp_template):
    """layout_index=99 exceeds what a default Presentation has — ManipulationError raised, no file written."""
    manifest = Manifest(
        name="test",
        type="pptx",
        template_file="template.pptx",
        description="test manifest",
        slide_types={
            "bad_slide": SlideTypeDef(
                layout_index=99,
                fields=[
                    FieldDef(name="title", placeholder_idx=0, kind="text", required=True),
                ],
            )
        },
    )
    content = {"slides": [{"type": "bad_slide", "fields": {"title": "Hi"}}]}
    out_path = tmp_path / "out.pptx"

    with pytest.raises(ManipulationError) as exc_info:
        render(manifest, content, out_path, tmp_template)

    assert "layout_index 99" in str(exc_info.value)
    assert not out_path.exists()


def test_render_bad_placeholder_idx_raises_manipulation_error(tmp_path, tmp_template):
    """placeholder_idx=99 absent in template layout — ManipulationError raised, no file written."""
    manifest = Manifest(
        name="test",
        type="pptx",
        template_file="template.pptx",
        description="test manifest",
        slide_types={
            "bad_slide": SlideTypeDef(
                layout_index=0,
                fields=[
                    FieldDef(name="title", placeholder_idx=99, kind="text", required=True),
                ],
            )
        },
    )
    content = {"slides": [{"type": "bad_slide", "fields": {"title": "Hi"}}]}
    out_path = tmp_path / "out.pptx"

    with pytest.raises(ManipulationError) as exc_info:
        render(manifest, content, out_path, tmp_template)

    assert "placeholder idx 99" in str(exc_info.value)
    assert not out_path.exists()


def test_render_template_missing_pptx_fails_cleanly(tmp_path, tmp_template):
    """Template .pptx deleted before render — ValueError raised, no output file written."""
    (tmp_template / "template.pptx").unlink()
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}
    out_path = tmp_path / "out.pptx"

    with pytest.raises(ManipulationError) as exc_info:
        render(manifest, content, out_path, tmp_template)

    assert "Template file not found" in str(exc_info.value)
    assert not out_path.exists()


def test_render_save_failure_leaves_no_partial_file(tmp_path, tmp_template):
    """prs.save() raises OSError — ManipulationError raised, no out_path or .pptx.tmp file left."""
    mock_prs = MagicMock()
    mock_prs.slide_layouts.__len__ = MagicMock(return_value=100)  # pass layout-index guard (layout_index=0 < 100)
    mock_prs.save.side_effect = OSError("disk full")
    out_path = tmp_path / "out.pptx"
    tmp_file = out_path.with_suffix(".pptx.tmp")

    with patch("templates_engine.render_pptx._open_template", return_value=mock_prs):
        manifest = load_manifest(tmp_template / "manifest.yaml")
        content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}

        with pytest.raises(ManipulationError) as exc_info:
            render(manifest, content, out_path, tmp_template)

    assert "disk full" in str(exc_info.value)
    assert not out_path.exists()
    assert not tmp_file.exists()


def test_render_rename_failure_leaves_no_temp_file(tmp_path, tmp_template):
    """Atomic replace raises OSError — ManipulationError raised, no out_path or .pptx.tmp left."""
    manifest = load_manifest(tmp_template / "manifest.yaml")
    content = {"slides": [{"type": "title_slide", "fields": {"title": "Hi"}}]}
    out_path = tmp_path / "out.pptx"
    tmp_file = out_path.with_suffix(".pptx.tmp")

    with patch("pathlib.Path.replace", side_effect=OSError("replace failed")):
        with pytest.raises(ManipulationError):
            render(manifest, content, out_path, tmp_template)

    assert not out_path.exists()
    assert not tmp_file.exists()
