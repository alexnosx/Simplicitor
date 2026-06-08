# tests/templates_engine/test_builtin_templates.py
# Phase L: Built-in template manifest validation and render smoke tests.
#
# Tier 1 (manifest-only): always run; manifest.yaml files are committed without
#   the template.pptx binaries, so these pass as soon as step 2 lands.
# Tier 2 (render + discovery): skipif the relevant template.pptx is absent, so
#   the suite stays green until Alex commits the binaries.
from pathlib import Path

import pytest
from pptx import Presentation

from templates_engine.config import get_builtin_root, list_templates
from templates_engine.manifest import load_manifest, lint_manifest
from templates_engine.render_pptx import render

_BUILTIN = get_builtin_root()
_BP_DIR = _BUILTIN / "business_pitch"
_TO_DIR = _BUILTIN / "technical_overview"
_BP_PPTX = _BP_DIR / "template.pptx"
_TO_PPTX = _TO_DIR / "template.pptx"

_skip_bp = pytest.mark.skipif(
    not _BP_PPTX.exists(),
    reason="business_pitch/template.pptx not yet committed",
)
_skip_to = pytest.mark.skipif(
    not _TO_PPTX.exists(),
    reason="technical_overview/template.pptx not yet committed",
)
_skip_both = pytest.mark.skipif(
    not (_BP_PPTX.exists() and _TO_PPTX.exists()),
    reason="both builtin template.pptx files must be committed",
)


# ---------------------------------------------------------------------------
# Tier 1 -- manifest-only (no .pptx dependency)
# ---------------------------------------------------------------------------

def test_business_pitch_manifest_loads():
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    assert manifest.name == "business_pitch"
    assert manifest.type == "pptx"


def test_technical_overview_manifest_loads():
    manifest = load_manifest(_TO_DIR / "manifest.yaml")
    assert manifest.name == "technical_overview"
    assert manifest.type == "pptx"


def test_business_pitch_has_four_slide_types():
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    assert set(manifest.slide_types.keys()) == {"title", "agenda", "content", "closing"}


def test_technical_overview_has_three_slide_types():
    manifest = load_manifest(_TO_DIR / "manifest.yaml")
    assert set(manifest.slide_types.keys()) == {"title", "architecture", "bullets"}


def test_business_pitch_lint_clean():
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    assert lint_manifest(manifest) == []


def test_technical_overview_lint_clean():
    manifest = load_manifest(_TO_DIR / "manifest.yaml")
    assert lint_manifest(manifest) == []


def test_business_pitch_title_subtitle_optional():
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    title = manifest.slide_types["title"]
    subtitle = next(f for f in title.fields if f.name == "subtitle")
    assert subtitle.required is False
    assert subtitle.kind == "text"
    assert subtitle.max_chars == 120


def test_business_pitch_agenda_items_required_max_five():
    """Charts template restored the bullet-capable OBJECT placeholders, so the
    agenda slide type carries a proper bullets field again (max 5 items)."""
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    agenda = manifest.slide_types["agenda"]
    items = next(f for f in agenda.fields if f.name == "items")
    assert items.required is True
    assert items.kind == "bullets"
    assert items.max_items == 5


def test_business_pitch_closing_statement_required():
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    closing = manifest.slide_types["closing"]
    statement = next(f for f in closing.fields if f.name == "statement")
    assert statement.required is True
    assert statement.kind == "text"
    assert statement.max_chars == 160


# ---------------------------------------------------------------------------
# Tier 2 -- business_pitch render tests (skipif .pptx absent)
# ---------------------------------------------------------------------------

@_skip_bp
def test_business_pitch_title_slide_renders(tmp_path):
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    content = {"slides": [{"type": "title", "fields": {"title": "Test Title"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", _BP_DIR)
    title_idx = next(
        f.placeholder_idx
        for f in manifest.slide_types["title"].fields
        if f.name == "title"
    )
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[title_idx].text == "Test Title"


@_skip_bp
def test_business_pitch_all_four_slide_types_render(tmp_path):
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    content = {
        "slides": [
            {"type": "title",   "fields": {"title": "Title"}},
            {"type": "agenda",  "fields": {"heading": "Agenda", "items": ["One", "Two", "Three"]}},
            {"type": "content", "fields": {"heading": "Content", "bullets": ["A", "B"]}},
            {"type": "closing", "fields": {"heading": "Thank You", "statement": "Reach out anytime."}},
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", _BP_DIR)
    prs = Presentation(str(result["path"]))
    assert len(prs.slides) == 4


@_skip_bp
def test_business_pitch_agenda_items_exact_order_and_count(tmp_path):
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    items = ["First item", "Second item", "Third item"]
    content = {"slides": [{"type": "agenda", "fields": {"heading": "Agenda", "items": items}}]}
    result = render(manifest, content, tmp_path / "out.pptx", _BP_DIR)
    prs = Presentation(str(result["path"]))
    agenda_items_idx = next(
        f.placeholder_idx
        for f in manifest.slide_types["agenda"].fields
        if f.name == "items"
    )
    texts = [
        p.text
        for p in prs.slides[0].placeholders[agenda_items_idx].text_frame.paragraphs
        if p.text
    ]
    assert texts == items


@_skip_bp
def test_business_pitch_title_subtitle_absent_no_issues(tmp_path):
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    content = {"slides": [{"type": "title", "fields": {"title": "Just a title"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", _BP_DIR)
    assert result["issues"] == []


@_skip_bp
def test_business_pitch_title_subtitle_written(tmp_path):
    manifest = load_manifest(_BP_DIR / "manifest.yaml")
    subtitle_value = "A bold step forward"
    content = {
        "slides": [
            {
                "type": "title",
                "fields": {
                    "title": "Acme Platform",
                    "subtitle": subtitle_value,
                },
            }
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", _BP_DIR)
    subtitle_idx = next(
        f.placeholder_idx
        for f in manifest.slide_types["title"].fields
        if f.name == "subtitle"
    )
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[subtitle_idx].text == subtitle_value


# ---------------------------------------------------------------------------
# Tier 2 -- technical_overview render tests (skipif .pptx absent)
# ---------------------------------------------------------------------------

@_skip_to
def test_technical_overview_title_slide_renders(tmp_path):
    manifest = load_manifest(_TO_DIR / "manifest.yaml")
    content = {"slides": [{"type": "title", "fields": {"title": "Overview"}}]}
    result = render(manifest, content, tmp_path / "out.pptx", _TO_DIR)
    prs = Presentation(str(result["path"]))
    assert prs.slides[0].placeholders[0].text == "Overview"


@_skip_to
def test_technical_overview_all_three_slide_types_render(tmp_path):
    manifest = load_manifest(_TO_DIR / "manifest.yaml")
    content = {
        "slides": [
            {"type": "title",        "fields": {"title": "Overview"}},
            {"type": "architecture", "fields": {"heading": "Arch", "components": ["Auth: handles login", "DB: stores data"]}},
            {"type": "bullets",      "fields": {"heading": "Details", "bullets": ["Point one", "Point two"]}},
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", _TO_DIR)
    prs = Presentation(str(result["path"]))
    assert len(prs.slides) == 3


@_skip_to
def test_technical_overview_architecture_components_exact_order_and_count(tmp_path):
    manifest = load_manifest(_TO_DIR / "manifest.yaml")
    components = ["Auth: validates tokens", "API: routes requests", "DB: persists state"]
    content = {
        "slides": [
            {"type": "architecture", "fields": {"heading": "Components", "components": components}}
        ]
    }
    result = render(manifest, content, tmp_path / "out.pptx", _TO_DIR)
    prs = Presentation(str(result["path"]))
    components_idx = manifest.slide_types["architecture"].fields[1].placeholder_idx
    texts = [
        p.text
        for p in prs.slides[0].placeholders[components_idx].text_frame.paragraphs
        if p.text
    ]
    assert texts == components


# ---------------------------------------------------------------------------
# Tier 2 -- discovery (skipif both .pptx absent)
# ---------------------------------------------------------------------------

@_skip_both
def test_builtins_appear_in_list_templates():
    results = list_templates(builtin_root=_BUILTIN)
    builtin_names = {t["name"] for t in results if t["source"] == "builtin"}
    assert "business_pitch" in builtin_names
    assert "technical_overview" in builtin_names


@_skip_both
def test_builtin_available_flag_true():
    results = list_templates(builtin_root=_BUILTIN)
    assert any(t["source"] == "builtin" for t in results)
