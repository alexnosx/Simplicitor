# tests/templates_engine/test_config.py
# Phase G: Tests for the two-directory template system.
from pathlib import Path
from unittest.mock import patch

import pytest
import yaml
from pptx import Presentation

from app.services.file_manipulator import ManipulationError
from templates_engine.config import (
    DEFAULT_TEMPLATE_NAMES,
    MANIFEST_NAME,
    TEMPLATE_PPTX_NAME,
    ensure_default_templates,
    get_app_data_dir,
    get_builtin_root,
    import_template,
    list_library,
    list_templates,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_pptx(path: Path) -> Path:
    """Save a default Presentation (has usable layouts) and return the path."""
    prs = Presentation()
    prs.save(str(path))
    return path


def _make_template_dir(root: Path, name: str) -> Path:
    """Create a fake but structurally valid template folder under *root*."""
    folder = root / name
    folder.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.save(str(folder / TEMPLATE_PPTX_NAME))
    # Minimal manifest YAML
    manifest = {
        "name": name,
        "type": "pptx",
        "template_file": TEMPLATE_PPTX_NAME,
        "description": "test",
        "slide_types": {},
    }
    (folder / MANIFEST_NAME).write_text(yaml.dump(manifest), encoding="utf-8")
    return folder


# ---------------------------------------------------------------------------
# get_builtin_root
# ---------------------------------------------------------------------------

def test_builtin_root_returns_path():
    assert isinstance(get_builtin_root(), Path)


def test_builtin_root_points_inside_package():
    root = get_builtin_root()
    assert "templates_engine" in str(root)
    assert root.name == "builtin"


# ---------------------------------------------------------------------------
# _ensure_dir (directory creation shared by import/list paths)
# ---------------------------------------------------------------------------

def test_user_root_creates_dir_if_missing(tmp_path):
    target = tmp_path / "user_templates"
    assert not target.exists()
    from templates_engine.config import _ensure_dir
    _ensure_dir(target)
    assert target.is_dir()


def test_user_root_returns_existing_dir(tmp_path):
    target = tmp_path / "user_templates"
    target.mkdir()
    from templates_engine.config import _ensure_dir
    _ensure_dir(target)  # must not raise when dir already exists
    assert target.is_dir()


def test_user_root_unwritable_raises_manipulation_error(tmp_path):
    target = tmp_path / "user_templates"
    from templates_engine.config import _ensure_dir
    with patch("pathlib.Path.mkdir", side_effect=OSError("permission denied")):
        with pytest.raises(ManipulationError, match=r"Cannot create"):
            _ensure_dir(target)


# ---------------------------------------------------------------------------
# list_templates
# ---------------------------------------------------------------------------

def test_list_templates_returns_list(tmp_path):
    result = list_templates(user_root=tmp_path / "user", builtin_root=tmp_path / "builtin")
    assert isinstance(result, list)


def test_list_templates_empty_roots_returns_empty(tmp_path):
    result = list_templates(user_root=tmp_path / "user", builtin_root=tmp_path / "builtin")
    assert result == []


def test_list_templates_detects_builtin(tmp_path):
    broot = tmp_path / "builtin"
    _make_template_dir(broot, "my_template")
    result = list_templates(user_root=tmp_path / "user", builtin_root=broot)
    assert len(result) == 1
    assert result[0]["source"] == "builtin"
    assert result[0]["name"] == "my_template"


def test_list_templates_detects_user(tmp_path):
    uroot = tmp_path / "user"
    _make_template_dir(uroot, "my_template")
    result = list_templates(user_root=uroot, builtin_root=tmp_path / "builtin")
    assert len(result) == 1
    assert result[0]["source"] == "user"


def test_list_templates_merges_both_roots(tmp_path):
    broot = tmp_path / "builtin"
    uroot = tmp_path / "user"
    _make_template_dir(broot, "builtin_tpl")
    _make_template_dir(uroot, "user_tpl")
    result = list_templates(user_root=uroot, builtin_root=broot)
    assert len(result) == 2
    sources = {t["source"] for t in result}
    assert sources == {"builtin", "user"}


def test_list_templates_builtin_comes_before_user(tmp_path):
    broot = tmp_path / "builtin"
    uroot = tmp_path / "user"
    _make_template_dir(broot, "aaa")
    _make_template_dir(uroot, "zzz")
    result = list_templates(user_root=uroot, builtin_root=broot)
    assert result[0]["source"] == "builtin"
    assert result[1]["source"] == "user"


def test_list_templates_ignores_dir_missing_pptx(tmp_path):
    uroot = tmp_path / "user"
    folder = uroot / "incomplete"
    folder.mkdir(parents=True)
    (folder / MANIFEST_NAME).write_text("name: x\n", encoding="utf-8")
    # No template.pptx → skipped
    result = list_templates(user_root=uroot, builtin_root=tmp_path / "builtin")
    assert result == []


def test_list_templates_ignores_dir_missing_manifest(tmp_path):
    uroot = tmp_path / "user"
    folder = uroot / "incomplete"
    folder.mkdir(parents=True)
    Presentation().save(str(folder / TEMPLATE_PPTX_NAME))
    # No manifest.yaml → skipped
    result = list_templates(user_root=uroot, builtin_root=tmp_path / "builtin")
    assert result == []


def test_list_templates_each_entry_has_required_keys(tmp_path):
    uroot = tmp_path / "user"
    _make_template_dir(uroot, "tpl")
    result = list_templates(user_root=uroot, builtin_root=tmp_path / "builtin")
    entry = result[0]
    assert "name" in entry
    assert "source" in entry
    assert "path" in entry
    assert "manifest_path" in entry
    assert "template_pptx" in entry


# ---------------------------------------------------------------------------
# import_template — happy path (usable deck)
# ---------------------------------------------------------------------------

def test_import_usable_returns_ok_status(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    result = import_template(pptx, user_root=tmp_path / "user")
    assert result["status"] == "ok"


def test_import_usable_creates_template_folder(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    uroot = tmp_path / "user"
    result = import_template(pptx, user_root=uroot)
    assert result["path"].is_dir()


def test_import_usable_folder_contains_template_pptx(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    result = import_template(pptx, user_root=tmp_path / "user")
    assert (result["path"] / TEMPLATE_PPTX_NAME).is_file()


def test_import_usable_folder_contains_manifest(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    result = import_template(pptx, user_root=tmp_path / "user")
    assert (result["path"] / MANIFEST_NAME).is_file()


def test_import_usable_result_has_report(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    result = import_template(pptx, user_root=tmp_path / "user")
    assert isinstance(result.get("report"), str)
    assert len(result["report"]) > 0


def test_import_usable_name_is_slug_of_stem(tmp_path):
    pptx = _make_pptx(tmp_path / "My Template.pptx")
    result = import_template(pptx, user_root=tmp_path / "user")
    assert result["name"] == "my_template"


# ---------------------------------------------------------------------------
# import_template — unusable deck
# ---------------------------------------------------------------------------

def _unusable_scoring():
    return {"is_usable": False, "layouts": []}


def test_import_unusable_returns_hard_stop_status(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    with patch("templates_engine.breakdown.score_layouts", return_value=_unusable_scoring()):
        result = import_template(pptx, user_root=tmp_path / "user")
    assert result["status"] == "hard_stop"


def test_import_unusable_writes_nothing(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    uroot = tmp_path / "user"
    with patch("templates_engine.breakdown.score_layouts", return_value=_unusable_scoring()):
        import_template(pptx, user_root=uroot)
    # The user root directory should have no template subdirectories.
    if uroot.exists():
        assert list(uroot.iterdir()) == []


# ---------------------------------------------------------------------------
# import_template — collision
# ---------------------------------------------------------------------------

def test_import_collision_returns_exists_status(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    uroot = tmp_path / "user"
    import_template(pptx, user_root=uroot)  # first import succeeds
    result = import_template(pptx, user_root=uroot)  # second collides
    assert result["status"] == "exists"
    assert result["name"] == "source"


# ---------------------------------------------------------------------------
# import_template — failure leaves no debris
# ---------------------------------------------------------------------------

def test_import_failed_strip_leaves_no_debris(tmp_path):
    pptx = _make_pptx(tmp_path / "source.pptx")
    uroot = tmp_path / "user"
    with patch(
        "templates_engine.breakdown.strip_to_template",
        side_effect=ManipulationError("simulated disk full"),
    ):
        with pytest.raises(ManipulationError):
            import_template(pptx, user_root=uroot)
    # No template folder should remain.
    template_dir = uroot / "source"
    assert not template_dir.exists()


def test_import_failed_manifest_write_leaves_no_debris(tmp_path):
    """Simulate manifest serialisation failure; no template folder should survive."""
    pptx = _make_pptx(tmp_path / "source.pptx")
    uroot = tmp_path / "user"
    # Patch yaml.dump in the config module so only manifest serialisation fails,
    # not strip_to_template or any other write in the call chain.
    with patch("templates_engine.config.yaml.dump", side_effect=OSError("disk full")):
        with pytest.raises((ManipulationError, OSError)):
            import_template(pptx, user_root=uroot)
    template_dir = uroot / "source"
    assert not template_dir.exists()


# ---------------------------------------------------------------------------
# ensure_default_templates
# ---------------------------------------------------------------------------

def _make_builtin_defaults(broot: Path) -> None:
    for name in DEFAULT_TEMPLATE_NAMES:
        _make_template_dir(broot, name)


def test_ensure_seeds_both_defaults_into_empty_root(tmp_path, monkeypatch):
    broot = tmp_path / "builtin"
    _make_builtin_defaults(broot)
    monkeypatch.setattr("templates_engine.config.get_builtin_root", lambda: broot)
    troot = tmp_path / "templates"
    ensure_default_templates(troot)
    for name in DEFAULT_TEMPLATE_NAMES:
        assert (troot / name / TEMPLATE_PPTX_NAME).is_file()
        assert (troot / name / MANIFEST_NAME).is_file()


def test_ensure_restores_deleted_default(tmp_path, monkeypatch):
    import shutil
    broot = tmp_path / "builtin"
    _make_builtin_defaults(broot)
    monkeypatch.setattr("templates_engine.config.get_builtin_root", lambda: broot)
    troot = tmp_path / "templates"
    ensure_default_templates(troot)
    shutil.rmtree(troot / DEFAULT_TEMPLATE_NAMES[0])
    ensure_default_templates(troot)
    assert (troot / DEFAULT_TEMPLATE_NAMES[0] / MANIFEST_NAME).is_file()


def test_ensure_leaves_existing_default_untouched(tmp_path, monkeypatch):
    broot = tmp_path / "builtin"
    _make_builtin_defaults(broot)
    monkeypatch.setattr("templates_engine.config.get_builtin_root", lambda: broot)
    troot = tmp_path / "templates"
    ensure_default_templates(troot)
    marker = troot / DEFAULT_TEMPLATE_NAMES[0] / MANIFEST_NAME
    marker.write_text(
        "name: edited\ntype: pptx\ntemplate_file: template.pptx\n"
        "description: e\nslide_types: {}\n",
        encoding="utf-8",
    )
    ensure_default_templates(troot)  # must not overwrite an existing default
    assert "edited" in marker.read_text(encoding="utf-8")


def test_ensure_idempotent(tmp_path, monkeypatch):
    broot = tmp_path / "builtin"
    _make_builtin_defaults(broot)
    monkeypatch.setattr("templates_engine.config.get_builtin_root", lambda: broot)
    troot = tmp_path / "templates"
    ensure_default_templates(troot)
    ensure_default_templates(troot)
    names = sorted(p.name for p in troot.iterdir() if p.is_dir())
    assert names == sorted(DEFAULT_TEMPLATE_NAMES)


def test_ensure_tolerates_missing_source(tmp_path, monkeypatch):
    broot = tmp_path / "builtin"
    broot.mkdir()  # empty: no default sources
    monkeypatch.setattr("templates_engine.config.get_builtin_root", lambda: broot)
    troot = tmp_path / "templates"
    ensure_default_templates(troot)  # must not raise
    assert not (troot / DEFAULT_TEMPLATE_NAMES[0]).exists()


# ---------------------------------------------------------------------------
# list_library
# ---------------------------------------------------------------------------

def test_list_library_tags_default_and_user(tmp_path):
    troot = tmp_path / "templates"
    _make_template_dir(troot, DEFAULT_TEMPLATE_NAMES[0])
    _make_template_dir(troot, "my_upload")
    by_name = {t["name"]: t["source"] for t in list_library(troot)}
    assert by_name[DEFAULT_TEMPLATE_NAMES[0]] == "default"
    assert by_name["my_upload"] == "user"


def test_list_library_skips_invalid(tmp_path):
    troot = tmp_path / "templates"
    (troot / "incomplete").mkdir(parents=True)
    assert list_library(troot) == []


def test_list_library_absent_root_returns_empty(tmp_path):
    assert list_library(tmp_path / "nonexistent") == []


def test_list_library_entry_has_required_keys(tmp_path):
    troot = tmp_path / "templates"
    _make_template_dir(troot, "tpl")
    entry = list_library(troot)[0]
    for key in ("name", "source", "path", "manifest_path", "template_pptx"):
        assert key in entry


# ── get_app_data_dir (shared config-dir resolver) ─────────────────────────────

def test_get_app_data_dir_uses_appdata_when_set(monkeypatch, tmp_path):
    monkeypatch.setenv("APPDATA", str(tmp_path / "Roaming"))
    assert get_app_data_dir() == tmp_path / "Roaming" / "Simplicitor"


def test_get_app_data_dir_falls_back_to_home_when_appdata_unset(monkeypatch):
    monkeypatch.delenv("APPDATA", raising=False)
    assert get_app_data_dir() == Path.home() / ".simplicitor"
