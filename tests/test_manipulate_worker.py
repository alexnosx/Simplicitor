# tests/test_manipulate_worker.py
import tempfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from app.services.ollama_client import OllamaConnectionError, OllamaGenerationError
from app.workers.manipulate_worker import ManipulateWorker


def make_worker(
    file_path: str = str(Path(tempfile.gettempdir()) / "doc.txt"),
    prompt: str = "Make it shorter",
    model: str = "llama3",
    client=None,
    backup_dir: str = str(Path(tempfile.gettempdir()) / "backups"),
) -> ManipulateWorker:
    if client is None:
        client = MagicMock()
        client.generate.return_value = "Modified content"
    return ManipulateWorker(
        file_path=file_path,
        prompt=prompt,
        model=model,
        client=client,
        backup_dir=backup_dir,
    )


def test_manipulate_worker_has_signals():
    w = make_worker()
    assert hasattr(w, "started")
    assert hasattr(w, "progress")
    assert hasattr(w, "completed")
    assert hasattr(w, "failed")


def test_manipulate_worker_emits_completed_on_success(qtbot, tmp_path):
    txt = tmp_path / "notes.txt"
    txt.write_text("Hello world", encoding="utf-8")
    backup_dir = tmp_path / "backups"
    client = MagicMock()
    client.generate.return_value = "Hello modified world"

    worker = ManipulateWorker(str(txt), "Make it better", "llama3", client, str(backup_dir))
    with qtbot.waitSignal(worker.completed, timeout=5000) as blocker:
        worker.run()

    saved_path, backup_path = blocker.args
    assert Path(saved_path).read_text(encoding="utf-8") == "Hello modified world"
    assert Path(backup_path).exists()
    assert Path(backup_path).read_text(encoding="utf-8") == "Hello world"


def test_manipulate_worker_emits_started(qtbot, tmp_path):
    txt = tmp_path / "f.txt"
    txt.write_text("x")
    client = MagicMock()
    client.generate.return_value = "y"
    worker = ManipulateWorker(str(txt), "change", "llama3", client, str(tmp_path / "bk"))
    started = []
    worker.started.connect(lambda: started.append(True))
    with qtbot.waitSignal(worker.completed, timeout=5000):
        worker.run()
    assert started


def test_manipulate_worker_emits_progress(qtbot, tmp_path):
    txt = tmp_path / "f.txt"
    txt.write_text("x")
    client = MagicMock()
    client.generate.return_value = "y"
    worker = ManipulateWorker(str(txt), "change", "llama3", client, str(tmp_path / "bk"))
    messages = []
    worker.progress.connect(messages.append)
    with qtbot.waitSignal(worker.completed, timeout=5000):
        worker.run()
    assert len(messages) >= 2


def test_manipulate_worker_fails_on_ollama_error(qtbot, tmp_path):
    txt = tmp_path / "f.txt"
    txt.write_text("x")
    client = MagicMock()
    client.generate.side_effect = OllamaConnectionError("refused")
    worker = ManipulateWorker(str(txt), "change", "llama3", client, str(tmp_path / "bk"))
    with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
        worker.run()
    assert "AI" in blocker.args[0]


def test_manipulate_worker_fails_on_empty_file(qtbot, tmp_path):
    txt = tmp_path / "empty.txt"
    txt.write_text("")
    client = MagicMock()
    worker = ManipulateWorker(str(txt), "change", "llama3", client, str(tmp_path / "bk"))
    with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
        worker.run()
    assert "empty" in blocker.args[0].lower()
    client.generate.assert_not_called()


def test_manipulate_worker_fails_on_missing_prompt_file(qtbot, tmp_path):
    txt = tmp_path / "f.txt"
    txt.write_text("content")
    client = MagicMock()
    worker = ManipulateWorker(str(txt), "change", "llama3", client, str(tmp_path / "bk"))
    with patch("app.workers.manipulate_worker.PROMPTS_DIR", tmp_path / "nonexistent"):
        with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
            worker.run()
    assert "configuration" in blocker.args[0].lower()
    client.generate.assert_not_called()


def test_manipulate_worker_emits_failed_on_unreadable_file(qtbot, tmp_path):
    f = tmp_path / "bad.docx"
    f.write_bytes(b"not a real docx")
    client = MagicMock()
    worker = ManipulateWorker(str(f), "change", "llama3", client, str(tmp_path / "bk"))
    with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
        worker.run()
    assert blocker.args[0]  # non-empty error message


# ── Scope detection tests ─────────────────────────────────────────────────────

def _make_valid_pptx(path: Path) -> None:
    """Write a minimal valid .pptx file with one slide for testing."""
    from pptx import Presentation as _Prs
    from pptx.util import Inches, Pt
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title+content layout
    slide.placeholders[0].text = "Test Slide"
    slide.placeholders[1].text = "Some bullet content"
    prs.save(str(path))


def _make_valid_docx(path: Path) -> None:
    """Write a minimal valid .docx file for testing."""
    from docx import Document as _Doc
    doc = _Doc()
    doc.add_paragraph("Hello world")
    doc.save(str(path))


def test_scope_check_rejects_styling_prompt_for_pptx(qtbot, tmp_path):
    """Styling keyword in prompt + .pptx → failed with scope message; Ollama not called."""
    pptx_file = tmp_path / "deck.pptx"
    _make_valid_pptx(pptx_file)
    client = MagicMock()

    worker = ManipulateWorker(
        str(pptx_file), "change the theme color to blue", "llama3", client, str(tmp_path / "bk")
    )
    with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
        worker.run()

    msg = blocker.args[0]
    assert "cannot" in msg.lower()
    assert "not modified" in msg.lower()
    client.generate.assert_not_called()
    assert not (tmp_path / "bk").exists()  # no backup directory created


def test_scope_check_rejects_styling_prompt_for_docx(qtbot, tmp_path):
    """Styling keyword in prompt + .docx → same scope rejection."""
    docx_file = tmp_path / "report.docx"
    _make_valid_docx(docx_file)
    client = MagicMock()

    worker = ManipulateWorker(
        str(docx_file), "update the font to Arial and make it bold", "llama3", client, str(tmp_path / "bk")
    )
    with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
        worker.run()

    msg = blocker.args[0]
    assert "cannot" in msg.lower()
    client.generate.assert_not_called()


def test_scope_check_does_not_block_txt_with_styling_keyword(qtbot, tmp_path):
    """Styling keyword in prompt + .txt → scope check does NOT fire; Ollama IS called."""
    txt_file = tmp_path / "notes.txt"
    txt_file.write_text("Some notes about color theory.", encoding="utf-8")
    client = MagicMock()
    client.generate.return_value = "Modified notes about color theory."

    worker = ManipulateWorker(
        str(txt_file), "mention the color blue more", "llama3", client, str(tmp_path / "bk")
    )
    with qtbot.waitSignal(worker.completed, timeout=5000):
        worker.run()

    client.generate.assert_called_once()


def test_scope_check_does_not_block_pptx_with_safe_prompt(qtbot, tmp_path):
    """Non-styling prompt + .pptx → scope check passes; pipeline continues normally."""
    pptx_file = tmp_path / "deck.pptx"
    _make_valid_pptx(pptx_file)
    client = MagicMock()
    client.generate.return_value = "[Slide 1]\nNew Title\nBullet one"

    worker = ManipulateWorker(
        str(pptx_file), "rewrite the slide titles to be more concise", "llama3", client, str(tmp_path / "bk")
    )
    with qtbot.waitSignal(worker.completed, timeout=5000):
        worker.run()

    client.generate.assert_called_once()


def test_scope_check_allows_word_containing_keyword(qtbot, tmp_path):
    """Keyword embedded in a longer word ("Colorado", "lifestyle") must not trigger."""
    pptx_file = tmp_path / "deck.pptx"
    _make_valid_pptx(pptx_file)
    client = MagicMock()
    client.generate.return_value = "[Slide 1]\nNew Title\nBullet one"

    worker = ManipulateWorker(
        str(pptx_file),
        "add a slide about lifestyle trends in Colorado",
        "llama3", client, str(tmp_path / "bk"),
    )
    with qtbot.waitSignal(worker.completed, timeout=5000):
        worker.run()

    client.generate.assert_called_once()


def test_scope_check_still_rejects_plural_keyword(qtbot, tmp_path):
    """Plural forms of out-of-scope keywords stay rejected."""
    pptx_file = tmp_path / "deck.pptx"
    _make_valid_pptx(pptx_file)
    client = MagicMock()

    worker = ManipulateWorker(
        str(pptx_file), "change the colors on every slide", "llama3", client, str(tmp_path / "bk")
    )
    with qtbot.waitSignal(worker.failed, timeout=5000) as blocker:
        worker.run()

    assert "cannot" in blocker.args[0].lower()
    client.generate.assert_not_called()


def test_scope_check_is_case_insensitive(qtbot, tmp_path):
    """Uppercase keywords are still caught."""
    pptx_file = tmp_path / "deck.pptx"
    _make_valid_pptx(pptx_file)
    client = MagicMock()

    worker = ManipulateWorker(
        str(pptx_file), "Change the THEME of this presentation", "llama3", client, str(tmp_path / "bk")
    )
    with qtbot.waitSignal(worker.failed, timeout=5000):
        worker.run()

    client.generate.assert_not_called()
