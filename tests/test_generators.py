# tests/test_generators.py
# Tests for Phase 3 Office document generators (Word, Excel, PowerPoint).
import pytest
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.generators.word_generator import WordGenerator
from app.generators.excel_generator import ExcelGenerator
from app.generators.pptx_generator import PptxGenerator


# ---------------------------------------------------------------------------
# WordGenerator tests
# ---------------------------------------------------------------------------

class TestWordGenerator:
    """Tests for WordGenerator.generate()."""

    def test_word_generator_creates_file(self, tmp_path: Path) -> None:
        parsed = {
            "title": "My Report",
            "sections": [{"heading": "Intro", "content": "Hello world.", "type": "text"}],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        assert out.exists()
        assert out.suffix == ".docx"

    def test_word_generator_returns_path_object(self, tmp_path: Path) -> None:
        parsed = {"title": "T", "sections": []}
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        assert isinstance(out, Path)

    def test_word_generator_title_as_heading1(self, tmp_path: Path) -> None:
        parsed = {"title": "Grand Title", "sections": []}
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        headings_1 = [p for p in doc.paragraphs if p.style.name == "Heading 1"]
        assert len(headings_1) == 1
        assert headings_1[0].text == "Grand Title"

    def test_word_generator_section_heading_as_heading2(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [{"heading": "Section A", "content": "body", "type": "text"}],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        headings_2 = [p for p in doc.paragraphs if "Heading 2" in p.style.name]
        assert any(h.text == "Section A" for h in headings_2)

    def test_word_generator_empty_heading_skipped(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [{"heading": "", "content": "text", "type": "text"}],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        headings = [p for p in doc.paragraphs if "Heading 2" in p.style.name]
        assert len(headings) == 0

    def test_word_generator_list_section(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [{"heading": "", "content": "Item 1\nItem 2", "type": "list"}],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        bullets = [p for p in doc.paragraphs if p.style.name == "List Bullet"]
        assert len(bullets) == 2

    def test_word_generator_list_skips_empty_lines(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [{"heading": "", "content": "Item 1\n\nItem 2\n\nItem 3", "type": "list"}],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        bullets = [p for p in doc.paragraphs if p.style.name == "List Bullet"]
        assert len(bullets) == 3

    def test_word_generator_text_section_paragraphs(self, tmp_path: Path) -> None:
        content = "First paragraph.\n\nSecond paragraph."
        parsed = {
            "title": "T",
            "sections": [{"heading": "", "content": content, "type": "text"}],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        normal_texts = [p.text for p in doc.paragraphs if p.style.name == "Normal"]
        assert "First paragraph." in normal_texts
        assert "Second paragraph." in normal_texts

    def test_word_generator_table_section(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [
                {
                    "heading": "Data",
                    "content": "Name\tAge\nAlice\t30\nBob\t25",
                    "type": "table",
                }
            ],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        assert len(doc.tables) == 1
        assert doc.tables[0].rows[0].cells[0].text == "Name"

    def test_word_generator_table_pipe_separated(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [
                {
                    "heading": "Data",
                    "content": "| Name | Age |\n| Alice | 30 |\n| Bob | 25 |",
                    "type": "table",
                }
            ],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        assert len(doc.tables) == 1
        assert doc.tables[0].rows[0].cells[0].text == "Name"

    def test_word_generator_table_header_bold(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "sections": [
                {
                    "heading": "D",
                    "content": "Name\tAge\nAlice\t30",
                    "type": "table",
                }
            ],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        header_cell = doc.tables[0].rows[0].cells[0]
        # At least one run in the header cell should be bold
        bold_runs = [r for r in header_cell.paragraphs[0].runs if r.bold]
        assert len(bold_runs) > 0

    def test_word_generator_creates_parent_dirs(self, tmp_path: Path) -> None:
        deep = tmp_path / "a" / "b" / "c" / "out.docx"
        out = WordGenerator().generate({"title": "T", "sections": []}, deep)
        assert out.exists()

    def test_word_generator_multiple_sections(self, tmp_path: Path) -> None:
        parsed = {
            "title": "Doc",
            "sections": [
                {"heading": "S1", "content": "Para 1.", "type": "text"},
                {"heading": "S2", "content": "Item A\nItem B", "type": "list"},
            ],
        }
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        headings_2 = [p for p in doc.paragraphs if "Heading 2" in p.style.name]
        assert len(headings_2) == 2
        bullets = [p for p in doc.paragraphs if p.style.name == "List Bullet"]
        assert len(bullets) == 2

    def test_word_generator_normal_style_font(self, tmp_path: Path) -> None:
        parsed = {"title": "T", "sections": []}
        out = WordGenerator().generate(parsed, tmp_path / "out.docx")
        doc = Document(str(out))
        normal_style = doc.styles["Normal"]
        assert normal_style.font.name == "Calibri"
        assert normal_style.font.size.pt == 11


# ---------------------------------------------------------------------------
# ExcelGenerator tests
# ---------------------------------------------------------------------------

class TestExcelGenerator:
    """Tests for ExcelGenerator.generate()."""

    def test_excel_generator_creates_file(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "Budget",
            "headers": ["Category", "Amount"],
            "rows": [["Food", "500"], ["Rent", "1200"]],
            "formulas": [],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        assert out.exists()
        assert out.suffix == ".xlsx"

    def test_excel_generator_returns_path_object(self, tmp_path: Path) -> None:
        parsed = {"sheet_name": "S", "headers": [], "rows": [], "formulas": []}
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        assert isinstance(out, Path)

    def test_excel_generator_sheet_name(self, tmp_path: Path) -> None:
        parsed = {"sheet_name": "My Sheet", "headers": ["A"], "rows": [], "formulas": []}
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        assert wb.active.title == "My Sheet"

    def test_excel_generator_headers_bold(self, tmp_path: Path) -> None:
        parsed = {"sheet_name": "S", "headers": ["A", "B"], "rows": [], "formulas": []}
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A1"].font.bold is True
        assert ws["B1"].font.bold is True

    def test_excel_generator_headers_in_row1(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["Name", "Score"],
            "rows": [],
            "formulas": [],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A1"].value == "Name"
        assert ws["B1"].value == "Score"

    def test_excel_generator_data_rows(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["X", "Y"],
            "rows": [["alpha", "beta"], ["gamma", "delta"]],
            "formulas": [],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A2"].value == "alpha"
        assert ws["B2"].value == "beta"
        assert ws["A3"].value == "gamma"

    def test_excel_generator_numeric_conversion(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["Val"],
            "rows": [["42"], ["3.14"]],
            "formulas": [],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A2"].value == 42
        assert ws["A3"].value == 3.14

    def test_excel_generator_non_numeric_stays_string(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["Label"],
            "rows": [["hello"], ["N/A"]],
            "formulas": [],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A2"].value == "hello"
        assert ws["A3"].value == "N/A"

    def test_excel_generator_formula_written(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["Value"],
            "rows": [["10"], ["20"]],
            "formulas": [{"cell": "A4", "formula": "=SUM(A2:A3)"}],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A4"].value == "=SUM(A2:A3)"

    def test_excel_generator_multiple_formulas(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["A", "B"],
            "rows": [["1", "2"]],
            "formulas": [
                {"cell": "A3", "formula": "=SUM(A2:A2)"},
                {"cell": "B3", "formula": "=SUM(B2:B2)"},
            ],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        assert ws["A3"].value == "=SUM(A2:A2)"
        assert ws["B3"].value == "=SUM(B2:B2)"

    def test_excel_generator_column_widths_set(self, tmp_path: Path) -> None:
        parsed = {
            "sheet_name": "S",
            "headers": ["LongHeaderName"],
            "rows": [["short"]],
            "formulas": [],
        }
        out = ExcelGenerator().generate(parsed, tmp_path / "out.xlsx")
        wb = load_workbook(str(out))
        ws = wb.active
        # Column A width should be set to at least len("LongHeaderName") + 2
        assert ws.column_dimensions["A"].width >= len("LongHeaderName") + 2

    def test_excel_generator_creates_parent_dirs(self, tmp_path: Path) -> None:
        deep = tmp_path / "x" / "y" / "out.xlsx"
        parsed = {"sheet_name": "S", "headers": [], "rows": [], "formulas": []}
        out = ExcelGenerator().generate(parsed, deep)
        assert out.exists()


# ---------------------------------------------------------------------------
# PptxGenerator tests
# ---------------------------------------------------------------------------

class TestPptxGenerator:
    """Tests for PptxGenerator.generate()."""

    def test_pptx_generator_creates_file(self, tmp_path: Path) -> None:
        parsed = {
            "title": "Deck",
            "slides": [{"title": "Title Slide", "bullets": [], "type": "title"}],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        assert out.exists()
        assert out.suffix == ".pptx"

    def test_pptx_generator_returns_path_object(self, tmp_path: Path) -> None:
        parsed = {"title": "T", "slides": []}
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        assert isinstance(out, Path)

    def test_pptx_generator_slide_count(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "slides": [
                {"title": "S1", "bullets": [], "type": "title"},
                {"title": "S2", "bullets": ["Point A", "Point B"], "type": "content"},
            ],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 2

    def test_pptx_generator_content_slide_has_bullets(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "slides": [{"title": "Content", "bullets": ["A", "B", "C"], "type": "content"}],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        texts.append(para.text)
        assert "A" in texts
        assert "B" in texts

    def test_pptx_generator_three_slide_types(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "slides": [
                {"title": "Title S", "bullets": ["subtitle"], "type": "title"},
                {"title": "Section S", "bullets": [], "type": "section"},
                {"title": "Content S", "bullets": ["Bullet"], "type": "content"},
            ],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 3

    def test_pptx_generator_empty_slides_list(self, tmp_path: Path) -> None:
        parsed = {"title": "Empty", "slides": []}
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        assert len(prs.slides) == 0

    def test_pptx_generator_content_slide_all_bullets(self, tmp_path: Path) -> None:
        bullets = ["First", "Second", "Third"]
        parsed = {
            "title": "T",
            "slides": [{"title": "S", "bullets": bullets, "type": "content"}],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        texts.append(para.text)
        for bullet in bullets:
            assert bullet in texts

    def test_pptx_generator_creates_parent_dirs(self, tmp_path: Path) -> None:
        deep = tmp_path / "p" / "q" / "out.pptx"
        parsed = {"title": "T", "slides": []}
        out = PptxGenerator().generate(parsed, deep)
        assert out.exists()

    def test_pptx_generator_title_slide_has_title_text(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "slides": [{"title": "My Title Slide", "bullets": [], "type": "title"}],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        title_texts = []
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                title_texts.append(shape.text_frame.text)
        assert any("My Title Slide" in t for t in title_texts)

    def test_pptx_generator_section_slide_has_title(self, tmp_path: Path) -> None:
        parsed = {
            "title": "T",
            "slides": [{"title": "Chapter 1", "bullets": [], "type": "section"}],
        }
        out = PptxGenerator().generate(parsed, tmp_path / "out.pptx")
        prs = Presentation(str(out))
        slide = prs.slides[0]
        title_texts = []
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                title_texts.append(shape.text_frame.text)
        assert any("Chapter 1" in t for t in title_texts)
