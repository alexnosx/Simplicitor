# tests/test_file_manipulator.py
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.config.defaults import MAX_MANIPULATION_CHARS
from app.services.file_manipulator import FileManipulator, ManipulationError


# ── Helpers ──────────────────────────────────────────────────────────────────

def make_docx(path: Path, text: str = "Hello world") -> Path:
    doc = Document()
    doc.add_paragraph(text)
    doc.save(str(path))
    return path


def make_xlsx(path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    wb.save(str(path))
    return path


def make_pptx(path: Path) -> Path:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.placeholders[0].text = "Title"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Bullet one\nBullet two"
    prs.save(str(path))
    return path


# ── extract_text ─────────────────────────────────────────────────────────────

def test_extract_txt(tmp_path):
    f = tmp_path / "notes.txt"
    f.write_text("Line one\nLine two", encoding="utf-8")
    result = FileManipulator().extract_text(f)
    assert "Line one" in result
    assert "Line two" in result


def test_extract_docx(tmp_path):
    f = make_docx(tmp_path / "doc.docx", "Sample paragraph text")
    result = FileManipulator().extract_text(f)
    assert "Sample paragraph text" in result


def test_extract_xlsx(tmp_path):
    f = make_xlsx(tmp_path / "data.xlsx")
    result = FileManipulator().extract_text(f)
    assert "Name" in result
    assert "Alice" in result


def test_extract_pptx(tmp_path):
    f = make_pptx(tmp_path / "deck.pptx")
    result = FileManipulator().extract_text(f)
    assert "Title" in result


def test_extract_unsupported_raises(tmp_path):
    f = tmp_path / "audio.mp3"
    f.write_bytes(b"fake")
    with pytest.raises(ManipulationError, match="Unsupported"):
        FileManipulator().extract_text(f)


def test_extract_truncates_large_file(tmp_path):
    f = tmp_path / "big.txt"
    f.write_text("A" * (MAX_MANIPULATION_CHARS + 5000), encoding="utf-8")
    result = FileManipulator().extract_text(f)
    assert len(result) == MAX_MANIPULATION_CHARS


def test_extract_corrupted_file_raises(tmp_path):
    f = tmp_path / "bad.docx"
    f.write_bytes(b"not a real docx file at all")
    with pytest.raises(ManipulationError):
        FileManipulator().extract_text(f)


# ── apply_changes ─────────────────────────────────────────────────────────────

def test_apply_txt_writes_modified_text(tmp_path):
    f = tmp_path / "notes.txt"
    f.write_text("original", encoding="utf-8")
    result = FileManipulator().apply_changes(f, "original", "modified content")
    assert result == f
    assert f.read_text(encoding="utf-8") == "modified content"


def test_apply_docx_creates_valid_docx(tmp_path):
    f = make_docx(tmp_path / "doc.docx")
    result = FileManipulator().apply_changes(f, "original", "New paragraph one\nNew paragraph two")
    assert result == f
    assert result.exists()
    doc = Document(str(result))
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    assert any("New paragraph" in t for t in texts)


def test_apply_pdf_saves_as_docx_alongside(tmp_path):
    f = tmp_path / "report.pdf"
    f.write_bytes(b"%PDF-1.4 fake")
    result = FileManipulator().apply_changes(f, "original text", "Updated content")
    assert result.suffix == ".docx"
    assert result.parent == tmp_path
    assert result.exists()


def test_apply_unsupported_raises(tmp_path):
    f = tmp_path / "audio.mp3"
    with pytest.raises(ManipulationError, match="Unsupported"):
        FileManipulator().apply_changes(f, "", "text")


def test_apply_xlsx_creates_valid_xlsx(tmp_path):
    f = make_xlsx(tmp_path / "data.xlsx")
    llm_response = "[Sheet: Results]\nName,Score\nBob,88"
    result = FileManipulator().apply_changes(f, "", llm_response)
    assert result == f
    # Just verify the file exists and is non-empty
    assert result.exists()
    assert result.stat().st_size > 0
