# simplicitor/app/services/file_manipulator.py
import logging
from pathlib import Path

from app.config.defaults import MAX_MANIPULATION_CHARS

logger = logging.getLogger(__name__)


class ManipulationError(Exception):
    """Raised when file text cannot be extracted or changes cannot be applied."""


class FileManipulator:
    """Extracts plain text from Office/text/PDF files and writes back LLM-modified content.

    Supports: .docx, .xlsx, .pptx, .txt, .pdf

    All public methods raise ManipulationError on failure so callers have a single
    exception type to handle.
    """

    # ── Public API ────────────────────────────────────────────────────────────

    def extract_text(self, file_path: Path) -> str:
        """Return plain text extracted from *file_path*.

        For .xlsx: comma-separated rows with [Sheet: name] headers.
        For .pptx: [Slide N] headers with title and bullet lines.
        For all others: plain paragraphs or lines.

        Truncates to MAX_MANIPULATION_CHARS if the file is very large.

        Args:
            file_path: Absolute path to the file.

        Returns:
            Extracted text string (may be empty for blank files).

        Raises:
            ManipulationError: If the file type is unsupported or the file
                cannot be read (corrupted, password-protected, etc.).
        """
        suffix = file_path.suffix.lower()
        extractors = {
            ".docx": self._extract_docx,
            ".xlsx": self._extract_xlsx,
            ".pptx": self._extract_pptx,
            ".txt": self._extract_txt,
            ".pdf": self._extract_pdf,
        }
        if suffix not in extractors:
            raise ManipulationError(f"Unsupported file type: {suffix!r}")
        try:
            return extractors[suffix](file_path)
        except ManipulationError:
            raise
        except Exception as exc:
            logger.debug("File extraction failed for %s: %s", file_path, exc)
            raise ManipulationError(f"Could not read file: {exc}") from exc

    def apply_changes(self, file_path: Path, original_text: str, llm_response: str) -> Path:
        """Write *llm_response* back to *file_path* (or a sibling for PDF).

        For .pdf: writes a .docx alongside the original PDF (never modifies PDF).

        Args:
            file_path: The file to overwrite with modified content.
            original_text: The original extracted text (unused by default
                handlers but available for future diffing).
            llm_response: The LLM-produced modified content.

        Returns:
            Path to the file that was written (same as file_path except PDF).

        Raises:
            ManipulationError: If the file type is unsupported or the write fails.
        """
        suffix = file_path.suffix.lower()
        handlers = {
            ".docx": self._apply_docx,
            ".xlsx": self._apply_xlsx,
            ".pptx": self._apply_pptx,
            ".txt": self._apply_txt,
            ".pdf": self._apply_pdf,
        }
        if suffix not in handlers:
            raise ManipulationError(f"Unsupported file type: {suffix!r}")
        try:
            return handlers[suffix](file_path, llm_response)
        except ManipulationError:
            raise
        except Exception as exc:
            logger.error("apply_changes failed for %s: %s", file_path, exc)
            raise ManipulationError(f"Could not apply changes: {exc}") from exc

    # ── Extractors ────────────────────────────────────────────────────────────

    def _extract_docx(self, path: Path) -> str:
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return self._truncate(text)

    def _extract_xlsx(self, path: Path) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            lines.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    lines.append(",".join(cells))
        return self._truncate("\n".join(lines))

    def _extract_pptx(self, path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        lines: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {idx}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return self._truncate("\n".join(lines))

    def _extract_txt(self, path: Path) -> str:
        text = path.read_text(encoding="utf-8", errors="replace")
        return self._truncate(text)

    def _extract_pdf(self, path: Path) -> str:
        import pdfplumber
        lines: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    lines.append(text)
        return self._truncate("\n".join(lines))

    # ── Writers ───────────────────────────────────────────────────────────────

    def _apply_docx(self, path: Path, text: str) -> Path:
        from docx import Document
        doc = Document()
        for para in text.strip().split("\n"):
            if para.strip():
                doc.add_paragraph(para.strip())
        doc.save(str(path))
        return path

    def _apply_xlsx(self, path: Path, text: str) -> Path:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        for line in text.strip().split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("[Sheet:") and line.endswith("]"):
                sheet_name = line[7:-1].strip()
                if ws.title == "Sheet1" and ws.max_row == 1 and ws.max_column == 1:
                    ws.title = sheet_name
                else:
                    ws = wb.create_sheet(title=sheet_name)
            else:
                cells = [c.strip() for c in line.split(",")]
                ws.append(cells)
        wb.save(str(path))
        return path

    def _apply_pptx(self, path: Path, text: str) -> Path:
        from pptx import Presentation
        from app.config.defaults import PPTX_LAYOUT_TITLE_CONTENT

        # Open the EXISTING file rather than Presentation() — a .pptx file is a
        # self-contained ZIP with its own slide masters and layouts, so no
        # default.pptx lookup occurs.  This also preserves the user's theme.
        prs = Presentation(str(path))

        # Remove all existing slides so we can rebuild from the LLM response.
        # drop_rel() cleans up the relationship entry; removing from _sldIdLst
        # removes the slide reference from the presentation manifest.
        slide_id_list = prs.slides._sldIdLst
        for sld_id in list(slide_id_list):
            r_id = sld_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if r_id:
                prs.part.drop_rel(r_id)
            slide_id_list.remove(sld_id)

        layout = prs.slide_layouts[PPTX_LAYOUT_TITLE_CONTENT]

        title = ""
        bullets: list[str] = []

        def flush() -> None:
            if not title and not bullets:
                return
            slide = prs.slides.add_slide(layout)
            if slide.placeholders:
                slide.placeholders[0].text = title
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for bullet in bullets:
                    tf.add_paragraph().text = bullet

        for line in text.strip().split("\n"):
            line = line.strip()
            if line.startswith("[Slide"):
                flush()
                title = ""
                bullets = []
            elif not title and line:
                title = line
            elif line:
                bullets.append(line)
        flush()

        prs.save(str(path))
        return path

    def _apply_txt(self, path: Path, text: str) -> Path:
        path.write_text(text, encoding="utf-8")
        return path

    def _apply_pdf(self, path: Path, text: str) -> Path:
        # Never modify a PDF; save the result as a .docx alongside it
        docx_path = path.with_suffix(".docx")
        return self._apply_docx(docx_path, text)

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _truncate(self, text: str) -> str:
        """Truncate text to MAX_MANIPULATION_CHARS if needed."""
        if len(text) > MAX_MANIPULATION_CHARS:
            logger.warning(
                "File content truncated from %d to %d chars", len(text), MAX_MANIPULATION_CHARS
            )
            return text[:MAX_MANIPULATION_CHARS]
        return text
