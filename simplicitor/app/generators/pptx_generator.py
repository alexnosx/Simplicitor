# simplicitor/app/generators/pptx_generator.py
# Phase 3: PowerPoint generator
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from app.config.defaults import (
    PPTX_LAYOUT_TITLE_SLIDE,
    PPTX_LAYOUT_TITLE_CONTENT,
    PPTX_LAYOUT_SECTION_HEADER,
)

logger = logging.getLogger(__name__)

# Bundled default template — avoids relying on pptx's internal package path,
# which is unreachable inside a Nuitka onefile executable.
_BUNDLED_TEMPLATE = Path(__file__).parent.parent.parent / "templates" / "pptx_default.pptx"


class PptxGenerator:
    """Generates .pptx files from parsed LLM output (Phase 3)."""

    def generate(self, parsed: dict, output_path: Path) -> Path:
        """Write a PowerPoint presentation from the parsed LLM response structure.

        Args:
            parsed: Validated dict from parse_pptx_response(). Expected keys:
                    ``title`` (str) and ``slides`` (list of slide dicts).
            output_path: Full file path (including filename) to write.

        Returns:
            output_path as a Path.

        Raises:
            OSError: On disk write failure.
        """
        output_path = Path(output_path)

        template_arg = str(_BUNDLED_TEMPLATE) if _BUNDLED_TEMPLATE.exists() else None
        prs = Presentation(template_arg)
        prs.core_properties.title = parsed.get("title", "")

        for idx, slide_data in enumerate(parsed.get("slides", [])):
            slide_type = slide_data.get("type", "content")
            slide_title = slide_data.get("title", "")
            # Normalise bullets: always a list of strings, never None
            bullets = slide_data.get("bullets") or []
            bullets = [str(b) for b in bullets if b is not None]

            try:
                if slide_type == "title":
                    self._add_title_slide(prs, slide_title, bullets)
                elif slide_type == "section":
                    self._add_section_slide(prs, slide_title)
                else:
                    # Default to content layout
                    self._add_content_slide(prs, slide_title, bullets)
            except Exception as exc:
                logger.error(
                    "Failed to add slide %d (type=%r, title=%r): %s", idx, slide_type, slide_title, exc
                )
                raise

        try:
            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
        except OSError as exc:
            logger.error("Failed to write %s: %s", output_path, exc)
            raise
        logger.debug("PowerPoint file written successfully: %s", output_path)
        return output_path

    def _add_title_slide(self, prs: Presentation, title: str, bullets: list[str]) -> None:
        """Add a title-slide (layout 0) with a title and optional subtitle."""
        layout = prs.slide_layouts[PPTX_LAYOUT_TITLE_SLIDE]
        slide = prs.slides.add_slide(layout)

        # Placeholder index 0 = title, index 1 = subtitle
        if slide.placeholders:
            title_ph = slide.placeholders[0]
            title_ph.text = title

        subtitle_text = bullets[0] if bullets else ""
        if len(slide.placeholders) > 1:
            subtitle_ph = slide.placeholders[1]
            subtitle_ph.text = subtitle_text

    def _add_section_slide(self, prs: Presentation, title: str) -> None:
        """Add a section-header slide (layout 2) with a title only."""
        layout = prs.slide_layouts[PPTX_LAYOUT_SECTION_HEADER]
        slide = prs.slides.add_slide(layout)

        if slide.placeholders:
            title_ph = slide.placeholders[0]
            title_ph.text = title

    def _add_content_slide(
        self, prs: Presentation, title: str, bullets: list[str]
    ) -> None:
        """Add a title-and-content slide (layout 1) with a title and bullet list."""
        layout = prs.slide_layouts[PPTX_LAYOUT_TITLE_CONTENT]
        slide = prs.slides.add_slide(layout)

        # Set title (placeholder idx 0)
        if slide.placeholders:
            slide.placeholders[0].text = title

        # Set bullet content (placeholder idx 1); skip if no bullets provided
        if len(slide.placeholders) > 1 and bullets:
            body_ph = slide.placeholders[1]
            tf = body_ph.text_frame
            tf.clear()

            for idx, bullet_text in enumerate(bullets):
                if idx == 0:
                    # Use the existing first paragraph
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.level = 0
                run = para.add_run()
                run.text = bullet_text
