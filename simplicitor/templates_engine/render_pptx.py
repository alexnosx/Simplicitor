# templates_engine/render_pptx.py
# Phase H: PPTX renderer.
import logging
import zipfile
from pathlib import Path

from templates_engine.manifest import Manifest, SlideTypeDef

logger = logging.getLogger(__name__)


def _open_template(path: Path):
    """Open a .pptx template, raising ValueError if missing, wrong type, or corrupt."""
    from pptx import Presentation
    from pptx.exceptions import InvalidXmlError, PackageNotFoundError

    if not path.exists():
        raise ValueError(f"Template file not found: '{path}'.")
    try:
        return Presentation(str(path))
    except (PackageNotFoundError, InvalidXmlError, zipfile.BadZipFile) as exc:
        raise ValueError(
            f"Could not open template '{path.name}' as a PowerPoint file."
        ) from exc
    except Exception as exc:
        raise ValueError(
            f"Could not open template '{path.name}' as a PowerPoint file "
            f"({type(exc).__name__})."
        ) from exc


def _get_placeholder(slide, idx: int, field_name: str, slide_idx: int):
    """Return the slide placeholder at *idx*, raising ManipulationError if absent."""
    from app.services.file_manipulator import ManipulationError

    try:
        return slide.placeholders[idx]
    except KeyError:
        raise ManipulationError(
            f"Slide {slide_idx}: placeholder idx {idx} (field '{field_name}') "
            f"not found in template. Manifest and template are out of sync."
        )


def _render_slide(
    slide,
    slide_def: SlideTypeDef,
    fields: dict,
    slide_idx: int,
) -> list[str]:
    """Populate slide placeholders from fields. Returns degrade warning strings."""
    issues: list[str] = []

    for field in slide_def.fields:
        idx = field.placeholder_idx
        name = field.name
        value = fields.get(name)

        if value is None:
            continue  # optional field absent — skip silently
        if field.kind == "bullets" and not value:
            continue  # empty bullet list — skip silently

        ph = _get_placeholder(slide, idx, name, slide_idx)

        if field.kind == "text":
            ph.text = value
            if field.max_chars is not None and len(value) > field.max_chars:
                msg = (
                    f"Slide {slide_idx}, field '{name}': text length {len(value)} "
                    f"exceeds max_chars {field.max_chars}."
                )
                logger.warning("%s", msg)
                issues.append(msg)

        elif field.kind == "bullets":
            tf = ph.text_frame
            tf.clear()
            for i, item in enumerate(value):
                if i == 0:
                    tf.paragraphs[0].text = item
                else:
                    tf.add_paragraph().text = item
            if field.max_items is not None and len(value) > field.max_items:
                msg = (
                    f"Slide {slide_idx}, field '{name}': {len(value)} bullets "
                    f"exceeds max_items {field.max_items}."
                )
                logger.warning("%s", msg)
                issues.append(msg)

        elif field.kind == "image":
            img_path = Path(value)
            if not img_path.exists():
                msg = (
                    f"Slide {slide_idx}, field '{name}': image path '{img_path}' "
                    f"not found, field skipped."
                )
                logger.warning("%s", msg)
                issues.append(msg)
                continue
            try:
                _ph = ph.insert_picture(str(img_path))  # noqa: F841 — capture; original ref invalidated
            except Exception as exc:
                msg = (
                    f"Slide {slide_idx}, field '{name}': could not insert image "
                    f"({type(exc).__name__}), field skipped."
                )
                logger.warning("%s", msg)
                issues.append(msg)

        else:
            msg = (
                f"Slide {slide_idx}, field '{name}': unknown kind '{field.kind}', "
                f"field skipped."
            )
            logger.warning("%s", msg)
            issues.append(msg)

    return issues


def render(
    manifest: Manifest,
    content: dict,
    out_path: str | Path,
    template_dir: str | Path,
) -> dict:
    """Render validated content into a PPTX deck using the named template.

    Content must already be validated via validate_content(). render() trusts
    the caller and does not re-validate.

    Args:
        manifest: Validated Manifest from load_manifest().
        content: Validated content dict {"slides": [{"type": str, "fields": dict}]}.
            Must be the parsed output of validate_content() — not raw JSON.
        out_path: Destination path. If suffix is empty, .pptx is appended silently.
        template_dir: Directory containing manifest.template_file.

    Returns:
        {"path": Path, "issues": list[str]} — path to the written file and
        any degrade warnings collected during rendering. Empty issues list
        means a clean render.

    Raises:
        ValueError: If the template .pptx is missing or corrupt.
        ManipulationError: If a layout_index or placeholder_idx from the
            manifest is absent in the template (manifest/template mismatch),
            or if the output file cannot be written. No partial file is left
            at out_path on failure.
    """
    from app.services.file_manipulator import ManipulationError

    out_path = Path(out_path)
    if not out_path.suffix:
        out_path = out_path.with_suffix(".pptx")
    template_dir = Path(template_dir)

    template_path = template_dir / manifest.template_file
    prs = _open_template(template_path)  # raises ValueError if missing/corrupt

    all_issues: list[str] = []

    for slide_idx, slide_data in enumerate(content["slides"]):
        slide_type = slide_data["type"]
        slide_def = manifest.slide_types[slide_type]
        layout_index = slide_def.layout_index

        if layout_index >= len(prs.slide_layouts):
            raise ManipulationError(
                f"Slide {slide_idx}: layout_index {layout_index} not found in template "
                f"(template has {len(prs.slide_layouts)} layout(s)). "
                f"Manifest and template are out of sync."
            )

        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)

        slide_issues = _render_slide(slide, slide_def, slide_data["fields"], slide_idx)
        all_issues.extend(slide_issues)

    # Save via temp file → atomic rename: guarantees no partial file at out_path on failure.
    tmp_file = out_path.with_suffix(".pptx.tmp")
    try:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(tmp_file))
        tmp_file.rename(out_path)
    except OSError as exc:
        try:
            tmp_file.unlink(missing_ok=True)
        except OSError:
            logger.warning("Could not clean up temp file '%s'.", tmp_file.name)
        logger.error("Failed to write rendered deck to '%s': %s", out_path, exc)
        raise ManipulationError(
            f"Could not write rendered deck to '{out_path.name}': {exc}"
        ) from exc

    logger.debug(
        "Rendered %d slide(s) to '%s' with %d issue(s).",
        len(content["slides"]),
        out_path.name,
        len(all_issues),
    )
    return {"path": out_path, "issues": all_issues}
