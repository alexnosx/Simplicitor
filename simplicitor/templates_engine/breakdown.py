# templates_engine/breakdown.py
# Phase D-F: PPTX structural inspector, content stripping, draft manifest generation.
import logging
import re
import zipfile
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

_EMU_PER_INCH = 914_400
_SLIDE_REL_ATTR = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
# Placeholder types that carry no fillable content (date stamp, footer text, slide counter).
_DECORATIVE_TYPES = frozenset({"DATE", "FOOTER", "SLIDE_NUMBER"})

_HARD_STOP_MESSAGE = (
    "This presentation can't be used as a template. Simplicitor builds slides by filling "
    "content placeholders defined in a presentation's layouts. This deck's slides use "
    "manually placed text boxes rather than layout placeholders, so there's no structure "
    "to fill. This is normal for hand-built decks - nothing is wrong with your file. You "
    "can use one of Simplicitor's built-in templates, or rebuild this deck using "
    "PowerPoint's slide layouts, then upload again."
)


def _emu_to_in(emu: int | None) -> str:
    if emu is None:
        return "?"
    return f"{emu / _EMU_PER_INCH:.2f}\""


def _slugify(text: str) -> str:
    """Convert a display name to a valid snake_case manifest key."""
    slug = re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_")
    return slug or "layout"


def _infer_kind(type_str: str) -> str:
    """Guess the field kind from a placeholder type string."""
    upper = type_str.upper()
    if "PICTURE" in upper or "MEDIA" in upper or "CLIP_ART" in upper:
        return "image"
    if "BODY" in upper:
        return "bullets"
    return "text"


def _label_placeholders(placeholders: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Assign name+kind to each non-decorative, non-custom placeholder in a layout.

    Labeling rules (in priority order):
    - idx=0 -> name "title", kind "text"
    - "PICTURE" in type -> name "image", kind "image"
    - Multiple placeholders sharing the same type string -> NEEDS_LABEL_<idx>
    - "BODY" in type (single) -> name "body", kind "bullets"
    - "SUBTITLE" in type (single) -> name "subtitle", kind "text"
    - Anything else (single, unrecognised) -> NEEDS_LABEL_<idx>, kind guessed

    Custom (idx>=10) and decorative (DATE/FOOTER/SLIDE_NUMBER) placeholders are skipped.
    """
    content_phs = [
        ph for ph in placeholders
        if ph["idx"] < 10
        and not any(dt in ph["type"].upper() for dt in _DECORATIVE_TYPES)
    ]

    # Count occurrences of each type string to detect ambiguity within this layout.
    type_counts: dict[str, int] = {}
    for ph in content_phs:
        type_counts[ph["type"]] = type_counts.get(ph["type"], 0) + 1

    fields: list[dict[str, Any]] = []
    for ph in content_phs:
        idx = ph["idx"]
        type_str = ph["type"]
        upper = type_str.upper()

        if idx == 0:
            name, kind = "title", "text"
        elif "PICTURE" in upper:
            name, kind = "image", "image"
        elif type_counts[type_str] > 1:
            name, kind = f"NEEDS_LABEL_{idx}", _infer_kind(type_str)
        elif "BODY" in upper:
            name, kind = "body", "bullets"
        elif "SUBTITLE" in upper:
            name, kind = "subtitle", "text"
        else:
            name, kind = f"NEEDS_LABEL_{idx}", _infer_kind(type_str)

        fields.append({
            "name": name,
            "placeholder_idx": idx,
            "kind": kind,
            "required": True,
        })

    return fields


def _open_presentation(path: Path):
    """Open a .pptx file, raising ValueError on missing, wrong extension, or corrupt file."""
    from pptx import Presentation
    from pptx.exceptions import InvalidXmlError, PackageNotFoundError

    if not path.exists():
        raise ValueError(f"File not found: '{path}'.")
    if path.suffix.lower() != ".pptx":
        raise ValueError(
            f"Expected a .pptx file, got '{path.suffix or '(no extension)'}' ('{path.name}')."
        )
    try:
        return Presentation(str(path))
    except (PackageNotFoundError, InvalidXmlError, KeyError, zipfile.BadZipFile) as exc:
        raise ValueError(
            f"Could not open '{path.name}' as a PowerPoint file."
        ) from exc
    except Exception as exc:
        # Unexpected failure — surface the type name only, never str(exc) which may
        # contain internal paths or unrelated details.
        raise ValueError(
            f"Could not open '{path.name}' as a PowerPoint file "
            f"({type(exc).__name__})."
        ) from exc


# ---------------------------------------------------------------------------
# Phase D: inspect_pptx + format_inspection
# ---------------------------------------------------------------------------

def inspect_pptx(path: str | Path) -> dict[str, Any]:
    """Inspect the layout/placeholder structure of a .pptx file (read-only).

    Args:
        path: Path to the .pptx file to inspect.

    Returns:
        A dict with keys:
            ``path`` (str): absolute resolved path.
            ``layouts`` (list): one entry per slide layout with:
                ``layout_index`` (int), ``name`` (str),
                ``placeholders`` (list) each with:
                    ``idx`` (int),
                    ``type`` (str, e.g. ``"TITLE (1)"``),
                    ``name`` (str),
                    ``position`` (dict: left/top/width/height in EMU, may be None),
                    ``is_custom`` (bool, True when idx >= 10).

    Raises:
        ValueError: If the file is missing, does not have a .pptx extension,
            or cannot be opened as a PowerPoint presentation.
    """
    path = Path(path)
    prs = _open_presentation(path)

    layouts: list[dict[str, Any]] = []
    for layout_idx, layout in enumerate(prs.slide_layouts):
        placeholders: list[dict[str, Any]] = []
        for ph in layout.placeholders:
            ph_idx = ph.placeholder_format.idx
            placeholders.append({
                "idx": ph_idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "position": {
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                },
                "is_custom": ph_idx >= 10,
            })
        layouts.append({
            "layout_index": layout_idx,
            "name": layout.name,
            "placeholders": placeholders,
        })

    logger.debug("Inspected '%s': %d layout(s).", path.name, len(layouts))
    return {
        "path": str(path.resolve()),
        "layouts": layouts,
    }


def format_inspection(report: dict[str, Any]) -> str:
    """Format an inspection report as a human-readable string for CLI output.

    Args:
        report: The dict returned by inspect_pptx().

    Returns:
        A multi-line string showing each layout and its placeholders.
    """
    lines: list[str] = [
        f"File: {report['path']}",
        f"Layouts: {len(report['layouts'])}",
    ]

    for layout in report["layouts"]:
        ph_count = len(layout["placeholders"])
        lines.append(
            f"\n  [{layout['layout_index']:2d}] {layout['name']}  "
            f"({ph_count} placeholder{'s' if ph_count != 1 else ''})"
        )
        for ph in layout["placeholders"]:
            pos = ph["position"]
            pos_str = (
                f"left={_emu_to_in(pos['left'])} "
                f"top={_emu_to_in(pos['top'])} "
                f"w={_emu_to_in(pos['width'])} "
                f"h={_emu_to_in(pos['height'])}"
            )
            custom = "  [CUSTOM]" if ph["is_custom"] else ""
            lines.append(
                f"       idx={ph['idx']:2d}  {ph['type']:<25}  {ph['name']!r:<35}  {pos_str}{custom}"
            )

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Phase E: strip_to_template, score_layouts
# ---------------------------------------------------------------------------

def strip_to_template(path: str | Path, out_path: str | Path) -> None:
    """Remove all slides from a .pptx, keeping masters/layouts/theme, and save.

    Args:
        path: Path to the source .pptx file.
        out_path: Destination path for the stripped template.

    Raises:
        ValueError: If path is missing, not a .pptx, or cannot be opened.
        ManipulationError: If out_path cannot be written (permissions, disk).
            No partial file is left behind on failure.
    """
    from app.services.file_manipulator import ManipulationError

    path = Path(path)
    out_path = Path(out_path)
    if out_path.suffix.lower() != ".pptx":
        raise ValueError(
            f"out_path must be a .pptx path, got '{out_path.suffix or '(no extension)'}'."
        )
    prs = _open_presentation(path)

    # Remove all slides while keeping masters/layouts/theme intact.
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        r_id = sld_id.get(_SLIDE_REL_ATTR)
        if r_id:
            prs.part.drop_rel(r_id)
        slide_id_list.remove(sld_id)

    try:
        prs.save(str(out_path))
    except OSError as exc:
        # Best-effort cleanup: remove any partial file written before the failure.
        try:
            out_path.unlink(missing_ok=True)
        except OSError:
            logger.warning("Could not clean up partial file '%s'.", out_path.name)
        logger.error("Failed to write stripped template to '%s': %s", out_path, exc)
        raise ManipulationError(
            f"Could not write template to '{out_path.name}': {exc}"
        ) from exc

    logger.debug("Stripped '%s' -> '%s' (0 slides).", path.name, out_path.name)


def score_layouts(inspection: dict[str, Any]) -> dict[str, Any]:
    """Score each layout in an inspection report as usable or unusable for templating.

    A layout is usable if it has at least one non-title content placeholder:
    idx >= 1 (not the title), idx < 10 (not custom), and type not in
    {DATE, FOOTER, SLIDE_NUMBER} (not purely decorative).
    Title-only layouts (only idx=0) are unusable — no body to fill.

    Args:
        inspection: The dict returned by inspect_pptx().

    Returns:
        A dict with:
            ``layouts`` (list): per-layout dicts with ``layout_index``, ``name``,
                ``usable`` (bool).
            ``is_usable`` (bool): True if at least one layout is usable.
    """
    if "layouts" not in inspection:
        raise ValueError("inspection must be the dict returned by inspect_pptx().")

    layout_scores: list[dict[str, Any]] = []
    any_usable = False

    for layout in inspection["layouts"]:
        has_content_ph = any(
            1 <= ph["idx"] < 10
            and not any(dt in ph["type"].upper() for dt in _DECORATIVE_TYPES)
            for ph in layout["placeholders"]
        )
        if has_content_ph:
            any_usable = True
        layout_scores.append({
            "layout_index": layout["layout_index"],
            "name": layout["name"],
            "usable": has_content_ph,
        })

    return {
        "layouts": layout_scores,
        "is_usable": any_usable,
    }


# ---------------------------------------------------------------------------
# Phase F: generate_draft_manifest, detection_report, hard_stop_result
# ---------------------------------------------------------------------------

def generate_draft_manifest(
    inspection: dict[str, Any],
    scoring: dict[str, Any],
    template_file: str,
) -> dict[str, Any]:
    """Generate a draft manifest dict from an inspection and scoring result.

    Only usable layouts (per scoring) are included as slide types.
    Placeholders are auto-labeled: idx=0 -> title, PICTURE -> image,
    same-type duplicates -> NEEDS_LABEL_<idx>. Custom and decorative
    placeholders are excluded. The returned dict is compatible with
    load_manifest after YAML serialisation.

    Args:
        inspection: The dict returned by inspect_pptx().
        scoring: The dict returned by score_layouts().
        template_file: Filename of the stripped .pptx (stored in the manifest).

    Returns:
        A manifest dict ready for yaml.dump / load_manifest round-trip.

    Raises:
        ValueError: If inspection or scoring are not the expected dicts.
    """
    if "layouts" not in inspection:
        raise ValueError("inspection must be the dict returned by inspect_pptx().")
    if "layouts" not in scoring:
        raise ValueError("scoring must be the dict returned by score_layouts().")

    usable_indices = {s["layout_index"] for s in scoring["layouts"] if s["usable"]}
    manifest_name = _slugify(Path(template_file).stem)

    slide_types: dict[str, Any] = {}
    slug_counts: dict[str, int] = {}

    for layout in inspection["layouts"]:
        if layout["layout_index"] not in usable_indices:
            continue
        slug = _slugify(layout["name"])
        count = slug_counts.get(slug, 0)
        slug_counts[slug] = count + 1
        key = slug if count == 0 else f"{slug}_{count}"
        slide_types[key] = {
            "layout_index": layout["layout_index"],
            "fields": _label_placeholders(layout["placeholders"]),
        }

    logger.debug(
        "Generated draft manifest for '%s': %d slide type(s).",
        template_file,
        len(slide_types),
    )
    return {
        "name": manifest_name,
        "type": "pptx",
        "template_file": template_file,
        "description": "Draft manifest — review NEEDS_LABEL fields before use.",
        "slide_types": slide_types,
    }


def detection_report(
    inspection: dict[str, Any],
    scoring: dict[str, Any],
) -> str:
    """Return a human-readable report of a deck's layout structure and usability.

    Args:
        inspection: The dict returned by inspect_pptx().
        scoring: The dict returned by score_layouts().

    Returns:
        A multi-line string suitable for CLI output or the GUI detection screen.

    Raises:
        ValueError: If inspection or scoring are not the expected dicts.
    """
    if "layouts" not in inspection:
        raise ValueError("inspection must be the dict returned by inspect_pptx().")
    if "layouts" not in scoring:
        raise ValueError("scoring must be the dict returned by score_layouts().")

    filename = Path(inspection.get("path", "(unknown)")).name
    usability_by_idx = {s["layout_index"]: s["usable"] for s in scoring["layouts"]}
    usable = [l for l in inspection["layouts"] if usability_by_idx.get(l["layout_index"])]
    unusable = [l for l in inspection["layouts"] if not usability_by_idx.get(l["layout_index"])]
    total = len(inspection["layouts"])

    lines: list[str] = [
        f"Inspection report: {filename}",
        f"Layouts: {total} total, {len(usable)} usable, {len(unusable)} unusable",
    ]

    if usable:
        # Pre-compute labels so _label_placeholders is called once per layout.
        lines.append("\nUsable layouts:")
        for layout in usable:
            fields = _label_placeholders(layout["placeholders"])
            field_summary = ", ".join(
                f"{f['name']} (idx={f['placeholder_idx']}, {f['kind']})" for f in fields
            )
            lines.append(
                f"  [{layout['layout_index']:2d}] {layout['name']}: "
                f"{field_summary or '(no fields)'}"
            )

    if unusable:
        lines.append("\nUnusable layouts (no fillable content placeholders):")
        for layout in unusable:
            lines.append(f"  [{layout['layout_index']:2d}] {layout['name']}")

    if scoring["is_usable"]:
        needs_label_count = sum(
            1
            for layout in usable
            for f in _label_placeholders(layout["placeholders"])
            if f["name"].startswith("NEEDS_LABEL_")
        )
        lines.append("\nVerdict: Deck CAN be used as a template.")
        if needs_label_count:
            lines.append(
                f"Action required: {needs_label_count} placeholder(s) need labelling "
                f"(marked NEEDS_LABEL_<idx>) — edit the manifest before use."
            )
    else:
        lines.append(
            "\nVerdict: Deck CANNOT be used as a template — "
            "no layout has fillable content placeholders."
        )

    return "\n".join(lines)


def hard_stop_result() -> dict[str, Any]:
    """Return the hard-stop result for a deck that cannot be used as a template.

    This is a normal returned value, not an exception — it represents expected
    user input (a hand-built deck with no layout placeholders), not a fault.
    The caller (import_template) decides when to return it.

    Returns:
        A dict with ``status`` "hard_stop" and the verbatim product-spec message.
    """
    return {
        "status": "hard_stop",
        "message": _HARD_STOP_MESSAGE,
    }
