# tests/templates_engine/test_pipeline.py
# Phase J: Tests for the generate-validate-repair-render pipeline.
import json
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation

from templates_engine.manifest import load_manifest
from templates_engine import pipeline

FIXTURE_MANIFEST = Path(__file__).parent / "fixtures" / "render_manifest.yaml"

# Minimal valid JSON that passes validate_content for render_manifest.yaml.
# title_slide requires: title (text, max_chars=20). "My Title" is 8 chars.
VALID_CONTENT = json.dumps({
    "slides": [{"type": "title_slide", "fields": {"title": "My Title"}}]
})

# Parseable JSON that fails validation: title_slide's required 'title' field is absent.
INVALID_CONTENT = json.dumps({
    "slides": [{"type": "title_slide", "fields": {}}]
})

# Structurally incomplete JSON: unbalanced braces → _looks_truncated returns True.
TRUNCATED_CONTENT = '{"slides": [{"type": "title_slide",'


@pytest.fixture
def manifest():
    return load_manifest(FIXTURE_MANIFEST)


def _assert_rendered_title(out_path: Path) -> None:
    """Open out_path and assert the first slide's title placeholder contains 'My Title'."""
    prs = Presentation(out_path)
    tf = prs.slides[0].placeholders[0].text_frame
    assert len(tf.paragraphs) == 1
    assert tf.paragraphs[0].text == "My Title"


# ---------------------------------------------------------------------------
# Happy path
# ---------------------------------------------------------------------------

def test_run_valid_first_attempt(manifest, tmp_template):
    """Pipeline succeeds on first attempt with no repair needed."""
    out_path = tmp_template / "output.pptx"
    with patch("templates_engine.llm.generate", return_value=VALID_CONTENT) as mock_gen:
        result = pipeline.run(manifest, tmp_template, [{"role": "user", "content": "deck"}], "llama3", out_path)

    assert mock_gen.call_count == 1
    assert result["path"] == out_path
    assert isinstance(result["issues"], list)
    _assert_rendered_title(out_path)


# ---------------------------------------------------------------------------
# Repair on parse failure
# ---------------------------------------------------------------------------

def test_run_parse_failure_repair_success(manifest, tmp_template):
    """Attempt 1 returns unparseable text; attempt 2 returns valid JSON. File is written."""
    out_path = tmp_template / "output.pptx"
    with patch("templates_engine.llm.generate", side_effect=["not json", VALID_CONTENT]) as mock_gen:
        result = pipeline.run(manifest, tmp_template, [{"role": "user", "content": "deck"}], "llama3", out_path)

    assert mock_gen.call_count == 2
    # Second call's messages must contain the parse-failure correction instruction.
    second_call_messages = mock_gen.call_args_list[1].args[0]
    correction_text = second_call_messages[-1]["content"]
    assert "could not be parsed" in correction_text
    _assert_rendered_title(out_path)


# ---------------------------------------------------------------------------
# Repair on validation failure
# ---------------------------------------------------------------------------

def test_run_validation_failure_repair_success(manifest, tmp_template):
    """Attempt 1 passes parse but fails validation; attempt 2 returns valid JSON."""
    out_path = tmp_template / "output.pptx"
    with patch("templates_engine.llm.generate", side_effect=[INVALID_CONTENT, VALID_CONTENT]) as mock_gen:
        result = pipeline.run(manifest, tmp_template, [{"role": "user", "content": "deck"}], "llama3", out_path)

    assert mock_gen.call_count == 2
    # Second call's messages must reference the field error (title is required).
    second_call_messages = mock_gen.call_args_list[1].args[0]
    full_correction = " ".join(m["content"] for m in second_call_messages)
    assert "title" in full_correction
    _assert_rendered_title(out_path)


# ---------------------------------------------------------------------------
# Hard fail after repair
# ---------------------------------------------------------------------------

def test_run_repair_still_fails_parse(manifest, tmp_template):
    """Both attempts return unparseable text → ParseError raised."""
    from app.parsers.llm_response_parser import ParseError

    out_path = tmp_template / "output.pptx"
    with patch("templates_engine.llm.generate", side_effect=["not json", "still not json"]):
        with pytest.raises(ParseError, match="after repair"):
            pipeline.run(manifest, tmp_template, [{"role": "user", "content": "deck"}], "llama3", out_path)

    assert not out_path.exists()


def test_run_repair_still_fails_validation(manifest, tmp_template):
    """Both attempts produce parseable JSON that fails validation → ParseError raised."""
    from app.parsers.llm_response_parser import ParseError

    out_path = tmp_template / "output.pptx"
    with patch("templates_engine.llm.generate", side_effect=[INVALID_CONTENT, INVALID_CONTENT]):
        with pytest.raises(ParseError, match="after repair"):
            pipeline.run(manifest, tmp_template, [{"role": "user", "content": "deck"}], "llama3", out_path)

    assert not out_path.exists()


# ---------------------------------------------------------------------------
# Truncation-bump
# ---------------------------------------------------------------------------

def test_run_truncation_bump_passes_max_tokens(manifest, tmp_template):
    """Truncated JSON on attempt 1 triggers max_tokens bump on the repair call."""
    from app.config.defaults import OLLAMA_REPAIR_MAX_TOKENS

    out_path = tmp_template / "output.pptx"
    with patch("templates_engine.llm.generate", side_effect=[TRUNCATED_CONTENT, VALID_CONTENT]) as mock_gen:
        pipeline.run(manifest, tmp_template, [{"role": "user", "content": "deck"}], "llama3", out_path)

    assert mock_gen.call_count == 2
    second_call_kwargs = mock_gen.call_args_list[1].kwargs
    assert "max_tokens" in second_call_kwargs, "Repair call must include max_tokens for truncation bump"
    assert second_call_kwargs["max_tokens"] >= OLLAMA_REPAIR_MAX_TOKENS
