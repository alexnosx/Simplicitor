"""Read-only PPTX placeholder inspector.

Usage:
    python scripts/inspect_template.py <path-to-template.pptx>

Prints, for each slide layout:
    - 0-based layout_index (use this value for layout_index in manifest.yaml)
    - layout name
    - each placeholder's idx, type, and shape name
      (use idx for placeholder_idx in manifest.yaml)

Run this after hand-building a template to reconcile manifest.yaml values
against what PowerPoint actually assigned. Makes no writes.
"""
import sys
from pathlib import Path


def _type_label(ph_type) -> str:
    name = getattr(ph_type, "name", None)
    if name:
        return f"{name}({int(ph_type)})"
    return str(ph_type)


def inspect(pptx_path: Path) -> None:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    print(f"Template : {pptx_path.name}")
    print(f"Layouts  : {len(prs.slide_layouts)}")
    print()
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}]  {layout.name!r}")
        placeholders = list(layout.placeholders)
        if placeholders:
            for ph in placeholders:
                fmt = ph.placeholder_format
                print(
                    f"        idx={fmt.idx}"
                    f"  type={_type_label(fmt.type)}"
                    f"  name={ph.name!r}"
                )
        else:
            print("        (no placeholders)")
        print()


def main() -> None:
    if len(sys.argv) != 2:
        print(__doc__)
        sys.exit(1)
    path = Path(sys.argv[1])
    if not path.exists():
        print(f"File not found: {path}")
        sys.exit(1)
    if path.suffix.lower() != ".pptx":
        print(f"Expected a .pptx file, got: {path.suffix!r}")
        sys.exit(1)
    inspect(path)


if __name__ == "__main__":
    main()
